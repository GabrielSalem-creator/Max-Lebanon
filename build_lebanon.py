# -*- coding: utf-8 -*-
import os, sys, io
from pathlib import Path

import subprocess
def _ensure(mod, pip=None):
    try:
        __import__(mod)
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", pip or mod, "-q"])
_ensure("pptx", "python-pptx")
_ensure("httpx")
_ensure("PIL", "pillow")

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent
# Load .env for TG token
for line in (ROOT / ".env").read_text().splitlines():
    line = line.strip()
    if line and not line.startswith("#") and "=" in line:
        k, v = line.split("=", 1)
        os.environ.setdefault(k.strip(), v.strip())
TG_TOKEN = os.environ["TG_TOKEN"]
chat_id = (ROOT / "data" / "tg_chat_id.txt").read_text().strip()

# Palette (Lebanese flag)
RED   = RGBColor(0xCE, 0x11, 0x26)
GREEN = RGBColor(0x00, 0x73, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY  = RGBColor(0x14, 0x1B, 0x2E)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
GOLD  = RGBColor(0xD4, 0xA0, 0x2A)

IMGS = {
    "beirut": "https://max.vdo-x.art/img/img_1781237650155.jpg",
    "cedars": "https://max.vdo-x.art/img/img_1781237652730.jpg",
    "food":   "https://max.vdo-x.art/img/img_1781237655447.jpg",
    "flag":   "https://max.vdo-x.art/img/img_1781237657905.jpg",
}
imgdir = ROOT / "data" / "_lebanon"
imgdir.mkdir(parents=True, exist_ok=True)
local = {}
with httpx.Client(timeout=30) as c:
    for k, u in IMGS.items():
        p = imgdir / f"{k}.jpg"
        p.write_bytes(c.get(u).content)
        local[k] = str(p)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, text, size, color, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb

def pic_cover(slide, path, x, y, w, h):
    # crop-to-fill into the box
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ratio = w / h
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic

def bullets(slide, x, y, w, h, items, size=20, color=NAVY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = "▸  " + head
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
        if body:
            r2 = p.add_run(); r2.text = "  " + body
            r2.font.size = Pt(size-2); r2.font.bold = False; r2.font.color.rgb = NAVY; r2.font.name = "Calibri"
    return tb

def header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.25), RED)
    rect(slide, 0, Inches(1.25), SW, Inches(0.08), WHITE)
    rect(slide, 0, Inches(1.33), SW, Inches(0.10), GREEN)
    txt(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.9), title, 32, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(slide, Inches(0.62), Inches(0.85), Inches(11), Inches(0.4), kicker, 13, RGBColor(0xFF,0xD9,0xDC), bold=False)
    # cedar emblem corner
    txt(slide, Inches(12.0), Inches(0.18), Inches(1.1), Inches(0.9), "\U0001F332", 30, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def footer(slide):
    rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), NAVY)
    txt(slide, Inches(0.5), SH - Inches(0.37), Inches(8), Inches(0.3),
        "LEBANON  •  Pearl of the Mediterranean", 11, WHITE, bold=False,
        anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
pic_cover(s, local["flag"], 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, NAVY).fill.fore_color.rgb  # overlay handled below
# semi overlay
ov = rect(s, 0, Inches(2.3), SW, Inches(3.2), NAVY)
ov.fill.fore_color.rgb = NAVY
ov.fill.transparency = 0  # solid band
txt(s, Inches(0.8), Inches(2.45), Inches(11.7), Inches(1.4), "LEBANON", 80, WHITE,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(5.16), Inches(3.95), Inches(3.0), Inches(0.06), RED)
txt(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.7), "Land of the Cedars  \U0001F332", 28, GOLD,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.5),
    "A journey through history, culture, nature & cuisine", 16, RGBColor(0xCF,0xD6,0xE4),
    bold=False, align=PP_ALIGN.CENTER)

# ---------- Slide 2: Quick Facts ----------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Quick Facts", "The essentials at a glance")
pic_cover(s, local["beirut"], Inches(7.6), Inches(1.7), Inches(5.2), Inches(4.9))
rect(s, Inches(7.6), Inches(1.7), Inches(0.12), Inches(4.9), RED)
bullets(s, Inches(0.65), Inches(1.85), Inches(6.7), Inches(5), [
    ("Capital", "Beirut — vibrant coastal capital"),
    ("Population", "~5.5 million people"),
    ("Official language", "Arabic (French & English widely spoken)"),
    ("Currency", "Lebanese Pound (LBP)"),
    ("Area", "10,452 km² — one of the smallest nations"),
    ("Government", "Parliamentary republic"),
], size=21)
footer(s)

# ---------- Slide 3: Geography & Nature ----------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Geography & Nature", "Mountains meeting the sea")
pic_cover(s, local["cedars"], Inches(0.55), Inches(1.7), Inches(5.2), Inches(4.9))
rect(s, Inches(5.63), Inches(1.7), Inches(0.12), Inches(4.9), GREEN)
bullets(s, Inches(6.1), Inches(1.85), Inches(6.7), Inches(5), [
    ("Mediterranean coast", "Sunny beaches along the western shore"),
    ("Mount Lebanon", "Snow-capped peaks — ski in the morning, swim by noon"),
    ("Cedars of God", "Ancient forests, the national symbol"),
    ("Bekaa Valley", "Fertile heartland & famous wine country"),
    ("Climate", "Mild Mediterranean, four real seasons"),
], size=21)
footer(s)

# ---------- Slide 4: History & Culture ----------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "History & Culture", "A 7,000-year-old crossroads")
rect(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(4.9), WHITE)
rect(s, Inches(0.55), Inches(1.7), Inches(0.12), Inches(4.9), RED)
bullets(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(4.5), [
    ("Phoenician roots", "Seafaring traders who gave the world the alphabet"),
    ("Byblos", "Among the oldest continuously inhabited cities on Earth"),
    ("Cultural mosaic", "18 official religious communities living side by side"),
    ("Arts & music", "Fairuz, poetry, cinema & a legendary nightlife"),
    ("French influence", "Architecture, cafes & a bilingual society"),
], size=22)
footer(s)

# ---------- Slide 5: Cuisine ----------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Cuisine", "One of the world's great food cultures")
pic_cover(s, local["food"], Inches(7.6), Inches(1.7), Inches(5.2), Inches(4.9))
rect(s, Inches(7.6), Inches(1.7), Inches(0.12), Inches(0.12), GREEN)
bullets(s, Inches(0.65), Inches(1.85), Inches(6.7), Inches(5), [
    ("Mezze", "Dozens of small shared plates"),
    ("Hummus & Tabbouleh", "Iconic dips & the fresh parsley salad"),
    ("Kibbeh", "National dish — spiced meat & bulgur"),
    ("Manakish", "Lebanon's beloved breakfast flatbread"),
    ("Bekaa wine & Arak", "Ancient winemaking & the anise spirit"),
    ("Sweets", "Baklava, knafeh & rosewater delights"),
], size=21)
footer(s)

# ---------- Slide 6: Landmarks & Tourism ----------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
header(s, "Landmarks & Tourism", "What to see")
rect(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(4.9), WHITE)
rect(s, Inches(0.55), Inches(1.7), Inches(0.12), Inches(4.9), GREEN)
bullets(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(4.5), [
    ("Baalbek", "Colossal Roman temples — a UNESCO wonder"),
    ("Jeita Grotto", "Breathtaking underground limestone caves"),
    ("Byblos", "Ancient harbor, crusader castle & souks"),
    ("Cedars of God", "Sacred mountain forest & ski resorts"),
    ("Beirut", "Corniche seafront, museums & nightlife"),
], size=22)
footer(s)

# ---------- Slide 7: Closing ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.5), RED)
rect(s, 0, SH - Inches(0.5), SW, Inches(0.5), GREEN)
txt(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6),
    "Lebanon", 70, WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8),
    "Resilient • Beautiful • Timeless", 26, GOLD, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.6),
    "\U0001F332  Land of the Cedars  \U0001F332", 22, RGBColor(0xCF,0xD6,0xE4), bold=False, align=PP_ALIGN.CENTER)

out = str(ROOT / "data" / "Lebanon.pptx")
prs.save(out)
print("SAVED:", out)

# ---------- Upload to Telegram ----------
with httpx.Client(timeout=60) as c:
    with open(out, "rb") as f:
        r = c.post(
            f"https://api.telegram.org/bot{TG_TOKEN}/sendDocument",
            data={"chat_id": chat_id, "caption": "\U0001F1F1\U0001F1E7 Lebanon — Pearl of the Mediterranean. Your presentation is ready!"},
            files={"document": ("Lebanon.pptx", f,
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
print("TG_STATUS:", r.status_code)
print("TG_RESP:", r.text[:300])
