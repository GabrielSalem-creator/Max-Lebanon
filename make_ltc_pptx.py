from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT2 = RGBColor(0xFF, 0x6B, 0x35)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xB0, 0xC4, 0xDE)
CARD    = RGBColor(0x1A, 0x2E, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_accent_line(slide, l, t, w, color=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_bullet_card(slide, bullets, l, t, w):
    row_h = 0.46
    total_h = len(bullets) * row_h + 0.3
    add_rect(slide, l, t, w, total_h, CARD)
    add_accent_line(slide, l, t, w, ACCENT)
    y = t + 0.18
    for b in bullets:
        dot = slide.shapes.add_shape(9, Inches(l+0.18), Inches(y+0.09), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        bx = slide.shapes.add_textbox(Inches(l+0.42), Inches(y), Inches(w-0.55), Inches(row_h))
        tf = bx.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        run = p2.add_run()
        run.text = b
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT
        y += row_h

slides_data = [
    {
        "type": "title",
        "title": "LTC6994",
        "subtitle": "TimerBlox Programmable Delay Block",
        "tagline": "A Clear University Guide — From Zero to Circuit-Ready"
    },
    {
        "type": "content",
        "num": "01",
        "title": "What Is the LTC6994?",
        "body": "A programmable delay chip from Analog Devices. It delays a digital signal by a time YOU set — using just resistors. No microcontroller. No code. Pure hardware precision.",
        "bullets": [
            "Part of the TimerBlox silicon timing IC family",
            "Delay range: 1 microsecond to 33.6 seconds",
            "Configured with only 1 to 3 external resistors",
            "Tiny SOT-23 or DFN-6 package — fits any board",
            "Supply voltage: 2.25V to 5.5V"
        ]
    },
    {
        "type": "content",
        "num": "02",
        "title": "How It Works — The Core Idea",
        "body": "Two internal blocks work together: a Master Oscillator and a Clock Divider. The oscillator ticks at a set frequency. The divider counts those ticks. When the count hits the target — the output fires.",
        "bullets": [
            "Step 1 — Input signal triggers the chip",
            "Step 2 — Internal oscillator starts ticking",
            "Step 3 — Clock divider (NDIV) scales the count",
            "Step 4 — Output fires after the full delay",
            "Result — Clean, accurate, delayed digital pulse"
        ]
    },
    {
        "type": "content",
        "num": "03",
        "title": "Key Pins — What Each One Does",
        "body": "Only 6 pins. Each one has a single clear job. Understand all 6 and you can wire this chip from scratch in 5 minutes.",
        "bullets": [
            "VDD  — Power supply input (2.25V to 5.5V)",
            "GND  — Ground reference",
            "RSET — Resistor that programs the master oscillator speed",
            "RDIV — Resistor that sets the clock divider value (NDIV)",
            "TRG  — Trigger input — this starts the delay countdown",
            "OUT  — The delayed output signal"
        ]
    },
    {
        "type": "content",
        "num": "04",
        "title": "Setting the Delay — The Formula",
        "body": "Delay equals NDIV divided by the oscillator frequency. RSET controls the oscillator. RDIV controls NDIV. Two resistors. One formula. Total control over timing.",
        "bullets": [
            "Formula:   Delay = NDIV / f_OSC",
            "RSET — higher resistance = slower oscillator = longer delay",
            "RDIV — programs NDIV: 8 settings from 1 to 2,097,152",
            "Small RSET + small NDIV = microsecond delays",
            "Large RSET + large NDIV = multi-second delays",
            "Datasheet Table 1 gives exact resistor values to use"
        ]
    },
    {
        "type": "content",
        "num": "05",
        "title": "LTC6994-1 vs LTC6994-2",
        "body": "Two versions. Same pinout. Same formula. Only the output polarity differs. Pick the version that matches your circuit logic.",
        "bullets": [
            "LTC6994-1 — Non-inverting: OUT follows TRG after the delay",
            "LTC6994-2 — Inverting: OUT is opposite of TRG after delay",
            "Both share identical pinout and resistor programming",
            "Use -1 for delay-then-activate circuits",
            "Use -2 for debouncing or delay-then-deactivate",
            "Choose based on required output logic polarity"
        ]
    },
    {
        "type": "content",
        "num": "06",
        "title": "Real-World Applications",
        "body": "The LTC6994 replaces bulky 555 timer circuits with a tiny, accurate, resistor-programmed chip. Found in almost every modern electronic product.",
        "bullets": [
            "Switch Debouncing — kills false button press triggers",
            "Power Sequencing — turns on voltage rails in safe order",
            "One-Shot Pulses — fires a single precise timed output",
            "Watchdog Timers — resets a system after timeout",
            "Motor Control — adds safe startup and shutdown delay",
            "Communication Systems — synchronizes signal timing"
        ]
    },
    {
        "type": "summary",
        "num": "07",
        "title": "Summary — 6 Points. Chip Mastered.",
        "bullets": [
            "1.   LTC6994 delays a digital signal from 1 microsecond to 33.6 seconds",
            "2.   Set the delay with 1 to 3 resistors — absolutely zero code needed",
            "3.   RSET controls the internal master oscillator frequency",
            "4.   RDIV programs the NDIV clock divider multiplier value",
            "5.   Version -1 is non-inverting, version -2 is inverting output",
            "6.   Used in debouncing, power sequencing, and one-shot timer circuits"
        ]
    }
]

blank_layout = prs.slide_layouts[6]

for sd in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG_DARK)

    if sd["type"] == "title":
        add_rect(slide, 0, 0, 0.45, 7.5, ACCENT)
        add_rect(slide, 0.45, 0, 0.18, 7.5, ACCENT2)
        add_text(slide, sd["title"], 1.0, 1.6, 11, 1.8, 80, bold=True, color=WHITE)
        add_text(slide, sd["subtitle"], 1.0, 3.5, 11, 0.7, 28, color=ACCENT)
        add_rect(slide, 1.0, 4.55, 11, 0.65, CARD)
        add_text(slide, sd["tagline"], 1.2, 4.62, 10.7, 0.52, 15, color=LIGHT)
        add_text(slide, "ANALOG DEVICES  |  TimerBlox  |  IC Reference", 1.0, 6.75, 11, 0.4, 11, color=ACCENT2)

    else:
        add_rect(slide, 0, 0, 13.33, 0.09, ACCENT)
        add_rect(slide, 0.3, 0.25, 0.72, 0.56, ACCENT)
        add_text(slide, sd["num"], 0.3, 0.25, 0.72, 0.56, 18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, sd["title"], 1.22, 0.2, 11.5, 0.65, 26, bold=True, color=WHITE)
        add_accent_line(slide, 1.22, 0.95, 11.0, ACCENT)

        if sd["type"] == "content":
            add_text(slide, sd["body"], 0.4, 1.08, 12.5, 0.88, 14, color=LIGHT)
            add_bullet_card(slide, sd["bullets"], 0.4, 2.05, 12.5)
        else:
            add_bullet_card(slide, sd["bullets"], 0.4, 1.12, 12.5)

        add_rect(slide, 0, 7.22, 13.33, 0.28, CARD)
        add_text(slide, "LTC6994  |  TimerBlox Delay Block  |  University Reference", 0.3, 7.22, 12.7, 0.28, 10, color=ACCENT2)

out = "C:/tmp/LTC6994_University_Guide.pptx"
os.makedirs("C:/tmp", exist_ok=True)
prs.save(out)
print("SAVED:", out)
