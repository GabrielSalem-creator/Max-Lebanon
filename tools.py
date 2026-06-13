#!/usr/bin/env python3
"""
MAX Tool Suite — call via: python3 /home/boxd/max/tools.py <tool> '<json>'
"""
import asyncio, json, os, re, sqlite3, subprocess, sys, uuid, tempfile as _tempfile
from datetime import datetime, timezone
from pathlib import Path

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

# Load .env from project root
_env_file = Path(__file__).parent / ".env"
if _env_file.exists():
    for _line in _env_file.read_text().splitlines():
        _line = _line.strip()
        if _line and not _line.startswith("#") and "=" in _line:
            _k, _v = _line.split("=", 1)
            os.environ.setdefault(_k.strip(), _v.strip())

# ── Credentials ───────────────────────────────────────────────────────────
RESEND_KEY = os.environ["RESEND_KEY"]
TG_TOKEN   = os.environ["TG_TOKEN"]
VID_URL    = os.environ.get("VID_URL", "https://veoaiapi.boxd.sh")
VID_KEY    = os.environ["VID_KEY"]
IMG_URL    = os.environ.get("IMG_URL", "https://image-z.created.app/api/generate-image")
LLM_URL    = os.environ.get("LLM_URL", "https://chat.good.hidns.vip/api/openai/v1")
LLM_MODEL  = os.environ.get("LLM_MODEL", "openai/gpt-oss-120b")
USER_EMAIL = os.environ.get("USER_EMAIL", "gsal4066@gmail.com")
IMAP_HOST  = os.environ.get("IMAP_HOST", "outlook.office365.com")
IMAP_EMAIL = os.environ["IMAP_EMAIL"]
IMAP_PASS  = os.environ.get("IMAP_PASS", "")

DATA_DIR   = Path(__file__).parent / "data"
TASKS_DB   = str(DATA_DIR / "tasks.db")


# ── Task Tracking ──────────────────────────────────────────────────────────
def _db():
    DATA_DIR.mkdir(exist_ok=True)
    conn = sqlite3.connect(TASKS_DB)
    conn.row_factory = sqlite3.Row
    conn.execute("""
        CREATE TABLE IF NOT EXISTS tasks (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            description TEXT DEFAULT '',
            status TEXT DEFAULT 'running',
            started_at TEXT NOT NULL,
            completed_at TEXT,
            duration_s REAL,
            result TEXT DEFAULT '',
            error TEXT DEFAULT ''
        )
    """)
    # Subtask hierarchy table for RSTD micro-tasking
    conn.execute("""
        CREATE TABLE IF NOT EXISTS subtasks (
            id TEXT PRIMARY KEY,
            parent_id TEXT NOT NULL,
            name TEXT NOT NULL,
            description TEXT DEFAULT '',
            tier INTEGER DEFAULT 1,
            status TEXT DEFAULT 'pending',
            validator TEXT DEFAULT '',
            retry_count INTEGER DEFAULT 0,
            max_retries INTEGER DEFAULT 3,
            result TEXT DEFAULT '',
            error TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """)
    conn.commit()
    return conn


async def task_create(p: dict) -> str:
    """
    Create a task record before starting work.
    p: {id: "slug", name: "short name", description: "what you're doing"}
    """
    tid = p.get("id") or str(uuid.uuid4())[:8]
    now = datetime.now(timezone.utc).isoformat()
    with _db() as conn:
        conn.execute(
            "INSERT OR REPLACE INTO tasks (id,name,description,status,started_at) VALUES (?,?,?,?,?)",
            (tid, p.get("name", "task"), p.get("description", ""), "running", now)
        )
    return f"✓ task:{tid} started at {now}"


async def task_update(p: dict) -> str:
    """
    Update task status after completion or failure.
    p: {id: str, status: "completed"|"failed"|"running", result: str, error: str}
    """
    tid = p.get("id", "")
    if not tid:
        return "✗ id required"
    status = p.get("status", "completed")
    now = datetime.now(timezone.utc).isoformat()
    with _db() as conn:
        row = conn.execute("SELECT started_at FROM tasks WHERE id=?", (tid,)).fetchone()
        duration = None
        if row:
            try:
                started = datetime.fromisoformat(row["started_at"])
                duration = (datetime.now(timezone.utc) - started).total_seconds()
            except Exception:
                pass
        conn.execute(
            "UPDATE tasks SET status=?,completed_at=?,duration_s=?,result=?,error=? WHERE id=?",
            (status,
             now if status != "running" else None,
             duration,
             p.get("result", ""),
             p.get("error", ""),
             tid)
        )
    dur_str = f" in {duration:.1f}s" if duration else ""
    return f"✓ task:{tid} → {status}{dur_str}"


async def task_list(p: dict) -> str:
    """
    List tasks. p: {status: "all"|"running"|"completed"|"failed", limit: 10}
    """
    status_filter = p.get("status", "all")
    limit = int(p.get("limit", 10))
    with _db() as conn:
        if status_filter == "all":
            rows = conn.execute(
                "SELECT * FROM tasks ORDER BY started_at DESC LIMIT ?", (limit,)
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT * FROM tasks WHERE status=? ORDER BY started_at DESC LIMIT ?",
                (status_filter, limit)
            ).fetchall()
    if not rows:
        return "No tasks found."
    lines = []
    for r in rows:
        dur = f"{r['duration_s']:.1f}s" if r['duration_s'] else "ongoing"
        preview = (r['result'] or r['error'] or "")[:120]
        lines.append(f"[{r['id']}] {r['name']} | {r['status']} | {dur}\n  {r['description']}\n  {preview}")
    return "\n\n".join(lines)


# ── Micro-Tasking (RSTD) ───────────────────────────────────────────────────

async def plan_task(p: dict) -> str:
    """
    Decompose a high-level goal into a subtask tree using the LLM.
    p: {parent_id: "slug", goal: "what to accomplish", context: "any relevant info"}
    Returns JSON subtask tree to execute sequentially or in parallel.
    """
    import httpx
    goal = p.get("goal", "")
    context = p.get("context", "")
    parent_id = p.get("parent_id", str(uuid.uuid4())[:8])

    prompt = f"""You are a task decomposition engine. Break this goal into the smallest possible atomic subtasks.

Goal: {goal}
Context: {context}

Rules:
- Each subtask must be independently verifiable (has a clear pass/fail check)
- Mark parallel=true for read-only or independent subtasks
- Mark parallel=false for tasks that mutate state or depend on prior results
- Assign tier: 1=simple scripted action, 2=complex reasoning needed, 3=human approval required
- Include a validator: what exact check confirms this subtask succeeded?

Respond with JSON only:
{{
  "parent_id": "{parent_id}",
  "goal": "{goal}",
  "subtasks": [
    {{"id":"s1","name":"short name","description":"what to do","tier":1,"parallel":false,"validator":"exact check to confirm success","depends_on":[]}}
  ]
}}"""

    async with httpx.AsyncClient(timeout=30) as c:
        r = await c.post(
            f"{LLM_URL}/chat/completions",
            headers={"Authorization": f"Bearer {VID_KEY}"},
            json={"model": LLM_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 1500}
        )
    if r.status_code != 200:
        return f"✗ plan_task LLM failed ({r.status_code}): {r.text[:200]}"

    try:
        content = r.json()["choices"][0]["message"]["content"].strip()
        # Extract JSON block
        import re as _re
        m = _re.search(r'\{[\s\S]*\}', content)
        if m:
            plan = json.loads(m.group())
            # Persist all subtasks as pending
            now = datetime.now(timezone.utc).isoformat()
            with _db() as conn:
                for st in plan.get("subtasks", []):
                    conn.execute(
                        "INSERT OR REPLACE INTO subtasks (id,parent_id,name,description,tier,status,validator,created_at,updated_at) VALUES (?,?,?,?,?,?,?,?,?)",
                        (st["id"], parent_id, st["name"], st.get("description",""), st.get("tier",1), "pending", st.get("validator",""), now, now)
                    )
            return json.dumps(plan, indent=2)
        return f"✗ plan_task: no JSON in response: {content[:200]}"
    except Exception as e:
        return f"✗ plan_task parse error: {e}"


async def subtask_create(p: dict) -> str:
    """
    Manually create a subtask under a parent task.
    p: {id, parent_id, name, description, tier (1-3), validator, max_retries}
    """
    sid = p.get("id") or str(uuid.uuid4())[:8]
    now = datetime.now(timezone.utc).isoformat()
    with _db() as conn:
        conn.execute(
            "INSERT OR REPLACE INTO subtasks (id,parent_id,name,description,tier,status,validator,max_retries,created_at,updated_at) VALUES (?,?,?,?,?,?,?,?,?,?)",
            (sid, p["parent_id"], p.get("name","subtask"), p.get("description",""),
             int(p.get("tier",1)), "pending", p.get("validator",""),
             int(p.get("max_retries",3)), now, now)
        )
    return f"✓ subtask:{sid} created under {p['parent_id']}"


async def subtask_update(p: dict) -> str:
    """
    Update subtask status after execution attempt.
    p: {id, status: pending|running|completed|failed|escalated, result, error}
    Status 'failed' auto-increments retry_count and sets status based on budget.
    """
    sid = p.get("id", "")
    if not sid:
        return "✗ id required"
    now = datetime.now(timezone.utc).isoformat()
    status = p.get("status", "completed")

    with _db() as conn:
        row = conn.execute("SELECT * FROM subtasks WHERE id=?", (sid,)).fetchone()
        if not row:
            return f"✗ subtask:{sid} not found"

        if status == "failed":
            new_retry = row["retry_count"] + 1
            if new_retry >= row["max_retries"]:
                tier = row["tier"]
                if tier < 3:
                    status = "escalated"
                    # Auto-create escalation record
                    conn.execute(
                        "UPDATE subtasks SET status=?,retry_count=?,error=?,updated_at=? WHERE id=?",
                        (status, new_retry, p.get("error",""), now, sid)
                    )
                    return f"✗ subtask:{sid} budget exhausted (tier {tier}) → ESCALATED to tier {tier+1}. Error: {p.get('error','')[:200]}"
                else:
                    status = "failed"
            conn.execute(
                "UPDATE subtasks SET status=?,retry_count=?,error=?,updated_at=? WHERE id=?",
                (status, new_retry, p.get("error",""), now, sid)
            )
            return f"✗ subtask:{sid} → {status} (retry {new_retry}/{row['max_retries']}). Error: {p.get('error','')[:200]}"

        conn.execute(
            "UPDATE subtasks SET status=?,result=?,error=?,updated_at=? WHERE id=?",
            (status, p.get("result",""), p.get("error",""), now, sid)
        )
    return f"✓ subtask:{sid} → {status}"


async def subtask_list(p: dict) -> str:
    """
    List subtasks for a parent task.
    p: {parent_id, status: "all"|"pending"|"running"|"completed"|"failed"|"escalated"}
    """
    parent_id = p.get("parent_id", "")
    status_filter = p.get("status", "all")
    with _db() as conn:
        if status_filter == "all":
            rows = conn.execute(
                "SELECT * FROM subtasks WHERE parent_id=? ORDER BY created_at ASC",
                (parent_id,)
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT * FROM subtasks WHERE parent_id=? AND status=? ORDER BY created_at ASC",
                (parent_id, status_filter)
            ).fetchall()
    if not rows:
        return f"No subtasks for {parent_id}"
    lines = []
    for r in rows:
        retry_info = f" (retry {r['retry_count']}/{r['max_retries']})" if r['retry_count'] > 0 else ""
        preview = (r['result'] or r['error'] or "")[:100]
        lines.append(f"[{r['id']}|T{r['tier']}] {r['name']} → {r['status']}{retry_info}\n  {r['description']}\n  validator: {r['validator']}\n  {preview}")
    return "\n\n".join(lines)


async def validate_result(p: dict) -> str:
    """
    Run a validation check on a subtask result. Returns pass/fail with detail.
    p: {subtask_id, result, validator_type: syntactic|semantic|tool|strategic, check: "what to verify"}

    Validator types:
    - syntactic: checks JSON/schema validity, regex match
    - semantic: asks LLM to verify result matches intent
    - tool: runs a bash command and checks exit code
    - strategic: checks if overall progress milestone reached
    """
    sid = p.get("subtask_id", "")
    result = p.get("result", "")
    vtype = p.get("validator_type", "semantic")
    check = p.get("check", "")

    if vtype == "syntactic":
        # Try JSON parse or regex
        try:
            json.loads(result)
            return f"✓ PASS syntactic: valid JSON"
        except Exception:
            if check:
                import re as _re
                if _re.search(check, result):
                    return f"✓ PASS syntactic: regex '{check}' matched"
                return f"✗ FAIL syntactic: regex '{check}' not found in: {result[:100]}"
            return f"✗ FAIL syntactic: invalid JSON and no regex check provided"

    elif vtype == "tool":
        # Run bash check command
        cmd = check
        try:
            proc = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=15)
            if proc.returncode == 0:
                return f"✓ PASS tool: exit 0\n{proc.stdout[:200]}"
            return f"✗ FAIL tool: exit {proc.returncode}\n{proc.stderr[:200]}"
        except Exception as e:
            return f"✗ FAIL tool error: {e}"

    elif vtype == "semantic":
        import httpx
        prompt = f"""Validation task. Does this result satisfy the requirement?

Requirement: {check}
Result: {result[:500]}

Answer with exactly: PASS or FAIL
Then one sentence explaining why."""
        async with httpx.AsyncClient(timeout=20) as c:
            r = await c.post(
                f"{LLM_URL}/chat/completions",
                headers={"Authorization": f"Bearer {VID_KEY}"},
                json={"model": LLM_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 100}
            )
        if r.status_code == 200:
            verdict = r.json()["choices"][0]["message"]["content"].strip()
            passed = verdict.upper().startswith("PASS")
            symbol = "✓" if passed else "✗"
            return f"{symbol} {'PASS' if passed else 'FAIL'} semantic: {verdict[:150]}"
        return f"✗ validator LLM error ({r.status_code})"

    elif vtype == "strategic":
        # Check if all subtasks for a parent are completed
        parent_id = p.get("parent_id", "")
        with _db() as conn:
            total = conn.execute("SELECT COUNT(*) as c FROM subtasks WHERE parent_id=?", (parent_id,)).fetchone()["c"]
            done = conn.execute("SELECT COUNT(*) as c FROM subtasks WHERE parent_id=? AND status='completed'", (parent_id,)).fetchone()["c"]
            failed = conn.execute("SELECT COUNT(*) as c FROM subtasks WHERE parent_id=? AND status IN ('failed','escalated')", (parent_id,)).fetchone()["c"]
        if failed > 0:
            return f"✗ FAIL strategic: {failed}/{total} subtasks failed"
        if done == total and total > 0:
            return f"✓ PASS strategic: {done}/{total} subtasks completed"
        return f"⏳ PARTIAL strategic: {done}/{total} completed, {total-done-failed} pending"

    return f"✗ unknown validator_type: {vtype}"


async def escalate_task(p: dict) -> str:
    """
    Escalate a stuck/failed subtask to the next resolution tier.
    p: {subtask_id, current_tier: 1|2|3, reason}
    Tier 1→2: retry with alt LLM model
    Tier 2→3: serialize state + notify user via Telegram
    """
    sid = p.get("subtask_id", "")
    tier = int(p.get("current_tier", 1))
    reason = p.get("reason", "task failed within retry budget")

    if tier == 1:
        # Escalate to T2: try with stronger model
        now = datetime.now(timezone.utc).isoformat()
        with _db() as conn:
            conn.execute(
                "UPDATE subtasks SET tier=2,status='pending',retry_count=0,updated_at=? WHERE id=?",
                (now, sid)
            )
        return f"↑ subtask:{sid} escalated T1→T2 (stronger model). Reset retry count. Retry now."

    elif tier == 2:
        # Escalate to T3: pause and notify user
        with _db() as conn:
            row = conn.execute("SELECT * FROM subtasks WHERE id=?", (sid,)).fetchone()
            conn.execute("UPDATE subtasks SET tier=3,status='escalated',updated_at=? WHERE id=?",
                         (datetime.now(timezone.utc).isoformat(), sid))

        state = {
            "subtask_id": sid,
            "name": row["name"] if row else sid,
            "description": row["description"] if row else "",
            "tier": 3,
            "reason": reason,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }
        # Notify user
        import httpx
        msg = f"MAX needs your help. Subtask '{state['name']}' is stuck after all retries. Reason: {reason}. Please review."
        async with httpx.AsyncClient(timeout=10) as c:
            chat_id_file = DATA_DIR / "tg_chat_id.txt"
            if chat_id_file.exists():
                chat_id = chat_id_file.read_text().strip()
                await c.post(
                    f"https://api.telegram.org/bot{TG_TOKEN}/sendMessage",
                    json={"chat_id": chat_id, "text": msg}
                )
        # Serialize state to disk for inspection
        state_file = DATA_DIR / f"escalation_{sid}.json"
        state_file.write_text(json.dumps(state, indent=2))
        return f"↑ subtask:{sid} escalated T2→T3 (human required). State saved to {state_file}. User notified via Telegram."

    return f"✗ already at tier 3 — human intervention required for subtask:{sid}"


# ── Email ──────────────────────────────────────────────────────────────────
async def send_email(p: dict) -> str:
    """Send email via Resend. Supports file attachments via 'attachment' (path) or 'attachments' (list of paths)."""
    import httpx, base64 as _b64
    payload = {
        "from": "MAX <onboarding@resend.dev>",
        "to": [p.get("to", USER_EMAIL)],
        "subject": p.get("subject", ""),
        "html": p.get("body", p.get("html", p.get("text", ""))),
    }
    # Support single attachment path or list
    attach_paths = []
    if p.get("attachment"):
        attach_paths.append(p["attachment"])
    if p.get("attachments"):
        attach_paths.extend(p["attachments"] if isinstance(p["attachments"], list) else [p["attachments"]])
    if attach_paths:
        attachments = []
        for path in attach_paths:
            try:
                fp = Path(path)
                if fp.exists():
                    content = _b64.b64encode(fp.read_bytes()).decode()
                    attachments.append({"filename": fp.name, "content": content})
            except Exception:
                pass
        if attachments:
            payload["attachments"] = attachments
    async with httpx.AsyncClient(timeout=30) as c:
        r = await c.post(
            "https://api.resend.com/emails",
            headers={"Authorization": f"Bearer {RESEND_KEY}"},
            json=payload,
        )
    if r.status_code in (200, 201):
        att_note = f" with {len(attach_paths)} attachment(s)" if attach_paths else ""
        return f"✓ Email sent to {p.get('to', USER_EMAIL)}{att_note} | id={r.json().get('id','ok')}"
    return f"✗ send_email failed ({r.status_code}): {r.text[:300]}"


async def read_emails(p: dict) -> str:
    count = p.get("count", 5)
    code = f"""
import imaplib, email
from email.header import decode_header
try:
    m = imaplib.IMAP4_SSL("{IMAP_HOST}", 993)
    m.login("{IMAP_EMAIL}", "{IMAP_PASS}")
    m.select("inbox")
    _, ids = m.search(None, "ALL")
    ids = ids[0].split()[-{count}:]
    out = []
    for uid in reversed(ids):
        _, d = m.fetch(uid, "(RFC822)")
        msg = email.message_from_bytes(d[0][1])
        subj = decode_header(msg["Subject"] or "")[0][0]
        if isinstance(subj, bytes): subj = subj.decode(errors="replace")
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True).decode(errors="replace")[:400]
                    break
        else:
            body = msg.get_payload(decode=True).decode(errors="replace")[:400]
        out.append(f"From: {{msg.get('From','')}}\\nDate: {{msg.get('Date','')}}\\nSubject: {{subj}}\\n{{body}}\\n---")
    m.close()
    print("\\n".join(out))
except Exception as e:
    print(f"Error: {{e}}")
"""
    r = subprocess.run([sys.executable, "-c", code], capture_output=True, text=True, timeout=30)
    return (r.stdout or r.stderr or "(no output)")[:3000]


# ── Telegram ───────────────────────────────────────────────────────────────
async def send_telegram_photo(p: dict) -> str:
    """
    Send a photo to Telegram. Accepts a public URL (max.vdo-x.art/img/...) or local file path.
    p: {url: "https://max.vdo-x.art/img/xxx.png", caption: "optional text"}
    Always use this after generate_image — sends the image inline, not just a link.
    """
    import httpx
    chat_file = DATA_DIR / "tg_chat_id.txt"
    chat_id = ""
    if chat_file.exists():
        chat_id = chat_file.read_text().strip()
    if not chat_id:
        async with httpx.AsyncClient(timeout=10) as c:
            r = await c.get(f"https://api.telegram.org/bot{TG_TOKEN}/getUpdates")
            updates = r.json().get("result", [])
            if updates:
                chat_id = str(updates[-1]["message"]["chat"]["id"])
                chat_file.write_text(chat_id)
    if not chat_id:
        return "✗ No Telegram chat_id — message the bot once first"

    photo_url = p.get("url", p.get("photo", ""))
    caption = p.get("caption", p.get("text", ""))

    # If it's a local file path, read and upload as multipart
    if photo_url and not photo_url.startswith("http"):
        local_path = Path(photo_url)
        if local_path.exists():
            async with httpx.AsyncClient(timeout=30) as c:
                with open(local_path, "rb") as f:
                    r = await c.post(
                        f"https://api.telegram.org/bot{TG_TOKEN}/sendPhoto",
                        data={"chat_id": chat_id, "caption": caption},
                        files={"photo": (local_path.name, f, "image/jpeg")},
                    )
            if r.status_code == 200:
                return f"✓ Photo sent to Telegram (file upload)"
            return f"✗ Telegram sendPhoto failed: {r.text[:200]}"

    # Send via URL (Telegram fetches it directly — works with public URLs)
    async with httpx.AsyncClient(timeout=15) as c:
        r = await c.post(
            f"https://api.telegram.org/bot{TG_TOKEN}/sendPhoto",
            json={"chat_id": chat_id, "photo": photo_url, "caption": caption},
        )
    if r.status_code == 200:
        return f"✓ Photo sent to Telegram: {photo_url}"
    return f"✗ Telegram sendPhoto failed ({r.status_code}): {r.text[:200]}"


async def send_telegram_document(p: dict) -> str:
    """
    Send a document/file (PPTX, PDF, DOCX, etc.) to Telegram.
    p: {path: "C:/tmp/file.pptx", caption: "optional text"}  (or 'file'/'document' key)
    Use this for presentations and files — NOT send_telegram_photo (which is images only).
    """
    import httpx
    chat_file = DATA_DIR / "tg_chat_id.txt"
    chat_id = ""
    if chat_file.exists():
        chat_id = chat_file.read_text().strip()
    if not chat_id:
        async with httpx.AsyncClient(timeout=10) as c:
            r = await c.get(f"https://api.telegram.org/bot{TG_TOKEN}/getUpdates")
            updates = r.json().get("result", [])
            if updates:
                chat_id = str(updates[-1]["message"]["chat"]["id"])
                chat_file.write_text(chat_id)
    if not chat_id:
        return "✗ No Telegram chat_id — message the bot once first"

    path = p.get("path", p.get("file", p.get("document", "")))
    caption = p.get("caption", p.get("text", ""))
    fp = Path(path)
    if not fp.exists():
        return f"✗ File not found: {path}"
    async with httpx.AsyncClient(timeout=60) as c:
        with open(fp, "rb") as f:
            r = await c.post(
                f"https://api.telegram.org/bot{TG_TOKEN}/sendDocument",
                data={"chat_id": chat_id, "caption": caption},
                files={"document": (fp.name, f)},
            )
    if r.status_code == 200:
        return f"✓ Document '{fp.name}' sent to Telegram"
    return f"✗ Telegram sendDocument failed ({r.status_code}): {r.text[:200]}"


async def send_whatsapp(p: dict) -> str:
    """
    Send a WhatsApp message via the logged-in WhatsApp Web (Chrome CDP on port 9222).
    p: {phone: "9613xxxxxx" (intl, no +), message: "text", file_url: "optional public link"}
    Run start_chrome_cdp.bat first so WhatsApp Web is logged in.
    Note: files are delivered as a link in the message (WhatsApp Web blocks scripted file attach).
    """
    import httpx, urllib.parse
    phone = str(p.get("phone", p.get("contact", p.get("number", "")))).lstrip("+").replace(" ", "")
    message = p.get("message", p.get("text", ""))
    if p.get("file_url"):
        message = (message + "\n" + p["file_url"]).strip()
    if not phone:
        return "✗ send_whatsapp needs a 'phone' (international format, e.g. 9613123456)"
    if not message:
        return "✗ send_whatsapp needs a 'message'"

    CDP = "http://localhost:9222"
    text_enc = urllib.parse.quote(message)
    url = f"https://web.whatsapp.com/send?phone={phone}&text={text_enc}"
    try:
        async with httpx.AsyncClient(timeout=10) as c:
            tabs = (await c.get(f"{CDP}/json")).json()
    except Exception as e:
        return f"✗ Chrome not reachable at {CDP}. Run start_chrome_cdp.bat first. Error: {e}"
    if not tabs:
        return "✗ No Chrome tabs found."

    import websockets
    # Prefer an existing WhatsApp tab, else use the first tab
    tab = next((t for t in tabs if "web.whatsapp.com" in t.get("url", "")), tabs[0])
    ws_url = tab.get("webSocketDebuggerUrl", "")
    if not ws_url:
        return "✗ No WebSocket debugger URL"
    mid = [1]
    async def send(ws, method, params=None):
        m = mid[0]; mid[0] += 1
        await ws.send(json.dumps({"id": m, "method": method, "params": params or {}}))
        while True:
            resp = json.loads(await asyncio.wait_for(ws.recv(), timeout=20))
            if resp.get("id") == m:
                return resp.get("result", {})
    try:
        async with websockets.connect(ws_url, max_size=10*1024*1024) as ws:
            await send(ws, "Page.navigate", {"url": url})
            # WhatsApp Web needs time to load the chat + prefilled text
            await asyncio.sleep(7)
            # Click the send button (data-icon="send")
            click_expr = (
                "(()=>{const b=document.querySelector('span[data-icon=\"send\"]')||"
                "document.querySelector('button[aria-label=\"Send\"]');"
                "if(b){b.click();return 'clicked';}return 'no-send-button';})()"
            )
            for _ in range(3):
                res = await send(ws, "Runtime.evaluate", {"expression": click_expr, "returnByValue": True})
                val = res.get("result", {}).get("value", "")
                if val == "clicked":
                    await asyncio.sleep(1)
                    return f"✓ WhatsApp message sent to {phone}"
                await asyncio.sleep(2)
            return ("⚠ Opened WhatsApp chat with the message prefilled for "
                    f"{phone}, but couldn't auto-click send — the chat is ready to send manually.")
    except Exception as e:
        return f"✗ send_whatsapp error: {e}"


async def send_telegram(p: dict) -> str:
    import httpx
    chat_id = p.get("chat_id", "")
    # Load saved chat_id from webhook capture
    chat_file = DATA_DIR / "tg_chat_id.txt"
    if not chat_id and chat_file.exists():
        chat_id = chat_file.read_text().strip()
    # Fallback: getUpdates
    if not chat_id:
        async with httpx.AsyncClient(timeout=10) as c:
            r = await c.get(f"https://api.telegram.org/bot{TG_TOKEN}/getUpdates")
            updates = r.json().get("result", [])
            if updates:
                chat_id = str(updates[-1]["message"]["chat"]["id"])
                chat_file.parent.mkdir(exist_ok=True)
                chat_file.write_text(chat_id)
    if not chat_id:
        return "✗ No chat_id — open Telegram, message @sickass_bot once, then retry"
    async with httpx.AsyncClient(timeout=10) as c:
        r = await c.post(
            f"https://api.telegram.org/bot{TG_TOKEN}/sendMessage",
            json={"chat_id": chat_id, "text": p["message"]}
        )
    if r.status_code == 200:
        return f"✓ Telegram sent to {chat_id}"
    return f"✗ Telegram failed: {r.text[:200]}"


PUBLIC_BASE = "https://max.vdo-x.art"


async def host_file(p: dict) -> str:
    """Copy a local file into the public web folder and return shareable links.
    For PPTX/Word/Excel it ALSO returns an Office Online viewer link that opens the
    presentation directly in any browser (phone or PC) — no app or download needed.
    p: {path: "C:/tmp/deck.pptx"}
    Returns the viewer link first (that's the one to send people).
    """
    import shutil, urllib.parse, time as _t
    src = Path(p.get("path", p.get("file", p.get("output", ""))))
    if not src.exists():
        return f"✗ File not found: {src}"
    files_dir = Path(__file__).parent / "static" / "files"
    files_dir.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[^A-Za-z0-9_.-]", "_", src.name)
    dest_name = f"{int(_t.time())}_{safe}"
    shutil.copy2(src, files_dir / dest_name)
    direct = f"{PUBLIC_BASE}/files/{dest_name}"
    ext = src.suffix.lower()
    if ext in (".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".xls"):
        viewer = "https://view.officeapps.live.com/op/view.aspx?src=" + urllib.parse.quote(direct, safe="")
        return json.dumps({"viewer_url": viewer, "direct_url": direct,
                           "message": f"✓ Hosted. Click to view in browser: {viewer}"})
    return json.dumps({"direct_url": direct, "message": f"✓ Hosted: {direct}"})


# ── Image Generation ───────────────────────────────────────────────────────
async def generate_image(p: dict) -> str:
    import httpx, base64, time
    async with httpx.AsyncClient(timeout=90) as c:
        r = await c.post(IMG_URL, json={"prompt": p["prompt"]})
    if r.status_code != 200:
        return f"✗ generate_image failed ({r.status_code}): {r.text[:200]}"
    d = r.json()
    raw = d.get("imageUrl") or d.get("url") or d.get("image_url") or ""
    if raw.startswith("data:image"):
        try:
            header, b64data = raw.split(",", 1)
            ext = "jpg" if "jpeg" in header else "png"
            fname = f"img_{int(time.time()*1000)}.{ext}"
            img_path = Path(__file__).parent / "static" / "img" / fname
            img_path.parent.mkdir(parents=True, exist_ok=True)
            img_path.write_bytes(base64.b64decode(b64data))
            public_url = f"{PUBLIC_BASE}/img/{fname}"
            return f"✓ Image generated: {public_url}"
        except Exception as e:
            return f"✗ generate_image: failed to save image: {e}"
    if raw:
        # External URL — download and re-host on max.vdo-x.art
        try:
            import time as _t
            async with httpx.AsyncClient(timeout=30) as c2:
                dl = await c2.get(raw)
            ct = dl.headers.get("content-type", "")
            ext = "jpg" if "jpeg" in ct else "png"
            fname = f"img_{int(_t.time()*1000)}.{ext}"
            img_path = Path(__file__).parent / "static" / "img" / fname
            img_path.parent.mkdir(parents=True, exist_ok=True)
            img_path.write_bytes(dl.content)
            public_url = f"{PUBLIC_BASE}/img/{fname}"
            return f"✓ Image generated: {public_url}"
        except Exception:
            return f"✓ Image generated: {raw}"
    return f"✗ generate_image: no URL in response: {str(d)[:200]}"


# ── Vondy Video Generation ────────────────────────────────────────────────
async def vondy_video(p: dict) -> str:
    """
    Free video gen via Vondy API with session rotation (no API key needed).
    p: {prompt: str, duration: "5"|"10", ratio: "16:9"|"9:16"|"1:1", resolution: "720p"|"1080p"}
    Downloads the video, saves to static/files/, returns public max.vdo-x.art URL.
    """
    import httpx, random as _rand, time as _t, re as _re

    prompt = p.get("prompt", "")
    duration = str(p.get("duration", "5"))
    ratio = p.get("ratio", p.get("aspectRatio", "16:9"))
    resolution = p.get("resolution", "720p")

    USER_AGENTS = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/147.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/146.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/147.0.0.0 Safari/537.36",
        "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/147.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:135.0) Gecko/20100101 Firefox/135.0",
    ]

    def _sec_ch_ua(ua):
        if "Firefox" in ua:
            return '"Firefox";v="135", "Not.A/Brand";v="8"'
        m = _re.search(r'Chrome/(\d+)', ua)
        v = m.group(1) if m else "147"
        return f'"Google Chrome";v="{v}", "Not.A/Brand";v="8", "Chromium";v="{v}"'

    def _platform(ua):
        return '"Windows"' if "Windows" in ua else '"macOS"' if "Mac" in ua else '"Linux"'

    def _device_id():
        return f"{uuid.uuid4().hex[:16]}-{uuid.uuid4().hex[:14]}-{format(int(_t.time()*1000),'x')}-{_rand.getrandbits(32):x}"

    def _extract_url(data):
        for key in ["url", "resultUrl", "videoUrl", "video_url"]:
            val = data.get(key, "")
            if isinstance(val, str) and val.startswith("http"):
                return val
        # nested data.url
        if isinstance(data.get("data"), dict):
            val = data["data"].get("url", "")
            if val and val.startswith("http"):
                return val
        # any .mp4 URL
        for v in data.values():
            if isinstance(v, str) and v.startswith("http") and ".mp4" in v:
                return v
        return None

    last_err = "no attempts made"
    for attempt in range(1, 6):
        ua = _rand.choice(USER_AGENTS)
        headers = {
            "Content-Type": "application/json",
            "User-Agent": ua,
            "sec-ch-ua": _sec_ch_ua(ua),
            "sec-ch-ua-mobile": "?0",
            "sec-ch-ua-platform": _platform(ua),
            "x-vondy-tier": "lite",
            "x-device-id": _device_id(),
            "Referer": "https://vondy.com/",
            "Origin": "https://vondy.com/",
            "Accept": "application/json, text/plain, */*",
            "Accept-Language": "en-GB,en-US;q=0.9,en;q=0.8",
        }
        payload = {
            "prompt": prompt, "duration": duration,
            "aspectRatio": ratio, "resolution": resolution, "cameraFixed": False,
        }
        try:
            async with httpx.AsyncClient(timeout=45) as c:
                r = await c.post("https://api.apivondy.com/api/open/video/generate",
                                 headers=headers, json=payload)
            if r.status_code == 429:
                last_err = f"rate limited (attempt {attempt})"
                await asyncio.sleep(attempt * 2)
                continue
            if not r.is_success:
                last_err = f"HTTP {r.status_code}: {r.text[:200]}"
                await asyncio.sleep(2)
                continue
            try:
                data = r.json()
            except Exception:
                last_err = f"bad JSON: {r.text[:100]}"
                continue
            video_url = _extract_url(data)
            if not video_url:
                last_err = f"no video URL in response: {str(data)[:200]}"
                continue

            # Download video and host on max.vdo-x.art
            try:
                async with httpx.AsyncClient(timeout=120, follow_redirects=True) as c:
                    dl = await c.get(video_url)
                fname = f"vondy_{uuid.uuid4().hex[:10]}.mp4"
                files_dir = Path(__file__).parent / "static" / "files"
                files_dir.mkdir(parents=True, exist_ok=True)
                (files_dir / fname).write_bytes(dl.content)
                public_url = f"{PUBLIC_BASE}/files/{fname}"
                return f"✓ vondy_video: {public_url}"
            except Exception as dl_err:
                return f"✓ vondy_video (direct URL): {video_url}\n(download failed: {dl_err})"

        except Exception as e:
            last_err = str(e)
            await asyncio.sleep(2)

    return f"✗ vondy_video failed after 5 attempts. Last error: {last_err}"


# ── Video Generation ───────────────────────────────────────────────────────
async def generate_video(p: dict) -> str:
    """Paid video gen via veoaiapi.boxd.sh (key sk-9661). p: {prompt, duration}"""
    import httpx
    async with httpx.AsyncClient(timeout=300) as c:
        r = await c.post(
            f"{VID_URL}/generate",
            headers={"Authorization": f"Bearer {VID_KEY}"},
            json={"prompt": p["prompt"], "duration": p.get("duration", 5)},
        )
    if r.status_code == 200:
        d = r.json()
        url = d.get("url") or d.get("video_url") or str(d)
        return f"✓ Video generated: {url}"
    return f"✗ generate_video failed: {r.text[:200]}"


async def generate_video_free(p: dict) -> str:
    """
    Free video gen via nanobananavideo.io. Runs in background by default.
    p: {prompt, size: "16:9"|"9:16"|"1:1", with_audio: bool, task_id: str, room: str}
    Returns immediately — background worker notifies when done via WS.
    """
    bg_script = Path(__file__).parent / "bg_video.py"
    params = json.dumps(p)
    proc = subprocess.Popen(
        [sys.executable, str(bg_script), params],
        stdout=open(os.path.join(_tempfile.gettempdir(), f"vidgen_{p.get('task_id','vid')}.log"), "w"),
        stderr=subprocess.STDOUT,
        start_new_session=True,
    )
    tid = p.get("task_id", "vid")
    return f"✓ Video generation started in background (task:{tid}, pid:{proc.pid}) — you'll be notified when ready"


# ── PowerPoint ─────────────────────────────────────────────────────────────
async def generate_powerpoint(p: dict) -> str:
    """
    Generate PPTX. p: {title, subtitle, slides:[{title, content, bullets:[]}], output}
    """
    script = f"""
import sys, subprocess
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches
import json

data = {json.dumps(p)}
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
sl = prs.slides.add_slide(prs.slide_layouts[0])
sl.shapes.title.text = data.get("title", "Presentation")
if len(sl.placeholders) > 1:
    sl.placeholders[1].text = data.get("subtitle", "")

# Content slides
for s in data.get("slides", []):
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = s.get("title", "")
    tf = sl.placeholders[1].text_frame
    tf.clear()
    if s.get("content"):
        tf.paragraphs[0].text = s["content"]
    for b in s.get("bullets", []):
        p2 = tf.add_paragraph()
        p2.text = str(b)
        p2.level = 1

out = data.get("output", "/tmp/max_presentation.pptx")
prs.save(out)
print(out)
"""
    r = subprocess.run([sys.executable, "-c", script], capture_output=True, text=True, timeout=60)
    if r.returncode == 0 and r.stdout.strip():
        return f"✓ PowerPoint created: {r.stdout.strip()}"
    return f"✗ generate_powerpoint failed: {r.stderr[:400]}"


# ── Web Scraper ────────────────────────────────────────────────────────────
async def scrape(p: dict) -> str:
    """
    Web scraper with two modes:
    - Query mode: p: {query: str, num_results: 5} — uses marsupilami-scraper API
    - URL mode:   p: {url: str, action: "text"|"html"} — fetches and parses page
    Falls back to DuckDuckGo search if marsupilami fails.
    """
    import httpx
    query = p.get("query", "")
    url = p.get("url", "")
    num = int(p.get("num_results", 5))

    if query:
        try:
            async with httpx.AsyncClient(timeout=30) as c:
                r = await c.post(
                    "https://v0-marsupilami-scraper.vercel.app/api/scrape",
                    headers={
                        "content-type": "application/json",
                        "sec-ch-ua": '"Chromium";v="116", "Not)A;Brand";v="24", "Google Chrome";v="116"',
                        "sec-ch-ua-mobile": "?0",
                        "sec-ch-ua-platform": '"macOS"',
                        "User-Agent": "Mozilla/5.0",
                        "origin": "https://v0-marsupilami-scraper.vercel.app",
                        "referer": "https://v0-marsupilami-scraper.vercel.app/",
                    },
                    json={"query": query, "num_results": num}
                )
            if r.status_code == 200:
                data = r.json()
                results = data.get("results", data if isinstance(data, list) else [])
                lines = []
                for item in (results if isinstance(results, list) else [])[:num]:
                    if isinstance(item, dict):
                        title = item.get("title", "")
                        snippet = item.get("snippet", item.get("description", ""))
                        link = item.get("link", item.get("url", ""))
                        content = str(item.get("content", ""))[:300]
                        parts = [x for x in [title, snippet, link, content] if x]
                        lines.append("\n  ".join(parts))
                    else:
                        lines.append(str(item)[:200])
                if lines:
                    return "\n\n".join(lines)
        except Exception:
            pass
        # Fallback to DuckDuckGo
        return await search({"query": query})

    if url:
        return await browse({"url": url, "action": p.get("action", "text")})

    return "✗ scrape: provide query or url"


# ── Web Search ─────────────────────────────────────────────────────────────
async def search(p: dict) -> str:
    import httpx
    query = p["query"]
    n = int(p.get("num_results", p.get("n", 8)))

    # Tier 1: SearXNG (self-hosted, unlimited, multi-engine)
    try:
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get("http://localhost:4040/search",
                            params={"q": query, "format": "json", "pageno": 1})
            if r.status_code == 200:
                data = r.json()
                results = data.get("results", [])[:n]
                if results:
                    lines = []
                    for x in results:
                        title   = x.get("title", "")
                        url     = x.get("url", "")
                        snippet = x.get("content", "")[:200]
                        lines.append(f"• {title}\n  {url}\n  {snippet}".strip())
                    return "\n\n".join(lines)
    except Exception:
        pass

    # Tier 2: DuckDuckGo HTML scraping
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as c:
            r = await c.get("https://html.duckduckgo.com/html/",
                            params={"q": query},
                            headers={"User-Agent": _UAS[0], "Accept": "text/html"})
        results = re.findall(
            r'class="result__a"[^>]*>(.*?)</a>.*?class="result__snippet"[^>]*>(.*?)</a>',
            r.text, re.DOTALL
        )
        if results:
            lines = []
            for title, snippet in results[:n]:
                t = re.sub(r"<[^>]+>", "", title).strip()
                s = re.sub(r"<[^>]+>", "", snippet).strip()
                lines.append(f"• {t}: {s}")
            return "\n".join(lines)
    except Exception:
        pass

    # Tier 3: Brave HTML scraping
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as c:
            r = await c.get("https://search.brave.com/search",
                            params={"q": query},
                            headers={"User-Agent": _UAS[1], "Accept": "text/html,*/*"})
        snippets = re.findall(r'<p class="snippet-description"[^>]*>(.*?)</p>', r.text, re.DOTALL)
        titles   = re.findall(r'<span class="snippet-title"[^>]*>(.*?)</span>', r.text, re.DOTALL)
        if snippets:
            lines = []
            for i, s in enumerate(snippets[:n]):
                t = re.sub(r"<[^>]+>", "", titles[i] if i < len(titles) else "").strip()
                s = re.sub(r"<[^>]+>", "", s).strip()
                lines.append(f"• {t}: {s}" if t else f"• {s}")
            return "\n".join(lines)
    except Exception:
        pass

    return f"✗ search: no results for '{query}'"


# ── Browse ─────────────────────────────────────────────────────────────────
_UAS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Mozilla/5.0 (iPhone; CPU iPhone OS 17_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.0 Mobile/15E148 Safari/604.1",
]

_HEADER_SETS = [
    {
        "User-Agent": _UAS[0],
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Sec-Fetch-Dest": "document",
        "Sec-Fetch-Mode": "navigate",
        "Sec-Fetch-Site": "none",
        "Upgrade-Insecure-Requests": "1",
    },
    {
        "User-Agent": _UAS[1],
        "Accept": "text/html,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.8",
        "Referer": "https://www.google.com/",
    },
    {
        "User-Agent": _UAS[3],
        "Accept": "text/html,application/xhtml+xml,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
    },
    {"User-Agent": _UAS[2]},
]


async def _fetch_robust(url: str) -> tuple[int, str]:
    """Try httpx with multiple header sets, then Playwright. Returns (status, html)."""
    import httpx
    for headers in _HEADER_SETS:
        try:
            async with httpx.AsyncClient(timeout=12, follow_redirects=True, verify=False) as c:
                r = await c.get(url, headers=headers)
                if r.status_code < 400:
                    return r.status_code, r.text
        except Exception:
            continue
    # Playwright fallback
    try:
        script = f"""
import asyncio
from playwright.async_api import async_playwright
async def run():
    async with async_playwright() as pw:
        b = await pw.chromium.launch(headless=True, args=["--no-sandbox"])
        p = await b.new_page()
        await p.goto({json.dumps(url)}, wait_until="domcontentloaded", timeout=20000)
        print(await p.content())
        await b.close()
asyncio.run(run())
"""
        r = subprocess.run([sys.executable, "-c", script], capture_output=True, text=True, timeout=35)
        if r.stdout.strip():
            return 200, r.stdout
    except Exception:
        pass
    return 0, ""


async def browse(p: dict) -> str:
    url = p["url"]
    action = p.get("action", "text")
    if action == "screenshot":
        out = os.path.join(_tempfile.gettempdir(), f"max_{abs(hash(url)) % 9999}.png")
        script = f"""
import asyncio
from playwright.async_api import async_playwright
async def run():
    async with async_playwright() as pw:
        b = await pw.chromium.launch(headless=True, args=["--no-sandbox"])
        p = await b.new_page()
        await p.goto({json.dumps(url)}, wait_until="domcontentloaded", timeout=20000)
        await p.screenshot(path={json.dumps(out)})
        await b.close()
asyncio.run(run())
"""
        subprocess.run([sys.executable, "-c", script], capture_output=True, timeout=30)
        return f"Screenshot: {out}"
    status, html = await _fetch_robust(url)
    if not html:
        return f"✗ browse: all methods failed for {url}"
    if action == "html":
        return html[:5000]
    text = re.sub(r"<script[^>]*>.*?</script>", "", html, flags=re.DOTALL)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:5000]


# ── Playwright Automation ──────────────────────────────────────────────────
async def automate(p: dict) -> str:
    url = p["url"]
    steps = p.get("steps", [])
    script = f"""
import asyncio, json
from playwright.async_api import async_playwright

async def run():
    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True, args=["--no-sandbox"])
        page = await browser.new_page()
        results = []
        try:
            await page.goto({json.dumps(url)}, wait_until="domcontentloaded", timeout=30000)
            results.append(f"Navigated to {json.dumps(url)}")
            for step in {json.dumps(steps)}:
                t = step.get("type", step.get("action",""))
                sel = step.get("selector", step.get("sel",""))
                val = step.get("value", step.get("val",""))
                ms  = step.get("ms", step.get("wait", 500))
                try:
                    if t == "click":
                        if sel.startswith("text:"):
                            await page.get_by_text(sel[5:]).first.click()
                        else:
                            await page.click(sel, timeout=8000)
                        results.append(f"clicked {{sel}}")
                    elif t == "fill":
                        if sel.startswith("label:"):
                            await page.get_by_label(sel[6:]).fill(val)
                        else:
                            await page.fill(sel, val)
                        results.append(f"filled {{sel}}")
                    elif t == "press":
                        await page.keyboard.press(val)
                        results.append(f"pressed {{val}}")
                    elif t == "select":
                        await page.select_option(sel, val)
                        results.append(f"selected {{val}}")
                    elif t == "wait":
                        await asyncio.sleep(ms/1000)
                        results.append(f"waited {{ms}}ms")
                    elif t in ("text","get_text"):
                        txt = await page.inner_text(sel or "body", timeout=5000)
                        results.append(f"text: {{txt[:400]}}")
                    elif t == "screenshot":
                        path = val or "/tmp/max_auto.png"
                        await page.screenshot(path=path)
                        results.append(f"screenshot: {{path}}")
                    elif t == "url":
                        results.append(f"url: {{page.url}}")
                    elif t == "scroll":
                        await page.evaluate(f"window.scrollBy(0, {{val or 500}})")
                        results.append(f"scrolled")
                except Exception as e:
                    results.append(f"err on {{t}} {{sel}}: {{e}}")
        except Exception as e:
            results.append(f"nav error: {{e}}")
        finally:
            await browser.close()
        print("\\n".join(results))

asyncio.run(run())
"""
    r = subprocess.run([sys.executable, "-c", script],
                       capture_output=True, text=True, timeout=90)
    out = r.stdout + (f"\n[err] {r.stderr[:300]}" if r.stderr.strip() else "")
    return out.strip()[:3000] or "(no output)"


# ── LLM (alt model, retries) ───────────────────────────────────────────────
async def llm(p: dict) -> str:
    """
    Call the free LLM API with retry on 504.
    p: {prompt: str, max_tokens: 800, model: "openai/gpt-oss-120b"}
    """
    import httpx, random
    headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
        "accept": "application/json",
        "content-type": "application/json",
        "origin": "https://chat.good.hidns.vip",
        "referer": "https://chat.good.hidns.vip/",
    }
    payload = {
        "messages": [{"role": "user", "content": p["prompt"]}],
        "model": p.get("model", LLM_MODEL),
        "stream": False,
        "max_tokens": p.get("max_tokens", 800),
    }
    async with httpx.AsyncClient(timeout=30, headers=headers) as c:
        for attempt in range(5):
            try:
                r = await c.post(f"{LLM_URL}/chat/completions", json=payload)
                if r.status_code == 200:
                    return r.json()["choices"][0]["message"]["content"]
                if r.status_code == 504:
                    wait = (2 ** attempt) + random.uniform(0, 1)
                    await asyncio.sleep(wait)
                    continue
                return f"✗ llm HTTP {r.status_code}: {r.text[:200]}"
            except Exception as e:
                if attempt < 4:
                    await asyncio.sleep(2 ** attempt)
                else:
                    return f"✗ llm error: {e}"
    return "✗ llm: all 5 retries failed"


# ── Web Auth & Session Management ──────────────────────────────────────────

SESSIONS_DIR = DATA_DIR / "sessions"
ACCOUNTS_FILE = DATA_DIR / "accounts.json"

# Stealth init script — removes webdriver flag, adds chrome object
_STEALTH_JS = """
Object.defineProperty(navigator, 'webdriver', {get: () => undefined});
window.chrome = {runtime: {}};
Object.defineProperty(navigator, 'plugins', {get: () => [1,2,3,4,5]});
Object.defineProperty(navigator, 'languages', {get: () => ['en-US','en']});
const getParameter = WebGLRenderingContext.prototype.getParameter;
WebGLRenderingContext.prototype.getParameter = function(parameter) {
  if (parameter === 37445) return 'Intel Inc.';
  if (parameter === 37446) return 'Intel Iris OpenGL Engine';
  return getParameter.call(this, parameter);
};
"""

import random as _random
import time as _time

def _jitter(base_ms=80, variance_ms=120):
    """Human-like delay: base + random variance."""
    return (base_ms + _random.random() * variance_ms) / 1000.0


async def _get_browser_context(session_name: str = None, headless: bool = True):
    """Launch a stealth Playwright browser context, optionally loading saved session."""
    from playwright.async_api import async_playwright
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)

    pw = await async_playwright().start()
    browser = await pw.chromium.launch(
        headless=headless,
        args=[
            "--no-sandbox",
            "--disable-blink-features=AutomationControlled",
            "--disable-dev-shm-usage",
        ]
    )

    session_file = SESSIONS_DIR / f"{session_name}.json" if session_name else None
    ctx_kwargs = {
        "user_agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
        "viewport": {"width": 1280, "height": 800},
        "locale": "en-US",
        "timezone_id": "America/New_York",
    }
    if session_file and session_file.exists():
        ctx_kwargs["storage_state"] = str(session_file)

    context = await browser.new_context(**ctx_kwargs)
    await context.add_init_script(_STEALTH_JS)
    return pw, browser, context


async def _save_session(context, session_name: str):
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    path = str(SESSIONS_DIR / f"{session_name}.json")
    await context.storage_state(path=path)
    return path


def _load_accounts() -> dict:
    if ACCOUNTS_FILE.exists():
        return json.loads(ACCOUNTS_FILE.read_text())
    return {}


def _save_accounts(data: dict):
    DATA_DIR.mkdir(exist_ok=True)
    ACCOUNTS_FILE.write_text(json.dumps(data, indent=2))


async def _smart_fill(page, field_type: str, value: str):
    """Fill a form field by trying common selectors + label text matching."""
    selectors = {
        "email":    ['input[type="email"]', 'input[name*="email"]', 'input[id*="email"]', 'input[placeholder*="email" i]'],
        "username": ['input[name*="user"]', 'input[id*="user"]', 'input[name*="login"]', 'input[placeholder*="username" i]'],
        "password": ['input[type="password"]'],
        "name":     ['input[name*="name"]', 'input[id*="name"]', 'input[placeholder*="name" i]'],
        "otp":      ['input[name*="otp"]', 'input[name*="code"]', 'input[name*="token"]', 'input[name*="verify"]', 'input[placeholder*="code" i]', 'input[maxlength="6"]'],
        "submit":   ['button[type="submit"]', 'input[type="submit"]', 'button:has-text("Sign in")', 'button:has-text("Log in")', 'button:has-text("Continue")', 'button:has-text("Next")', 'button:has-text("Sign up")', 'button:has-text("Create account")'],
    }
    for sel in selectors.get(field_type, []):
        try:
            el = page.locator(sel).first
            if await el.is_visible(timeout=2000):
                if field_type == "submit":
                    await el.click()
                    return True
                await el.click()
                await page.wait_for_timeout(int(_jitter() * 1000))
                await el.fill("")
                for char in value:
                    await el.type(char, delay=_random.randint(40, 120))
                return True
        except Exception:
            continue
    return False


async def web_login(p: dict) -> str:
    """
    Login to a website. Saves session on success for reuse.
    p: {url, username, password, session_name?, screenshot?}
    Returns: success/fail + session_name
    """
    url = p.get("url", "")
    username = p.get("username", "")
    password = p.get("password", "")
    session_name = p.get("session_name") or url.split("/")[2].replace("www.", "").replace(".", "_")

    if not url or not username or not password:
        return "✗ web_login requires: url, username, password"

    # Check for existing valid session
    session_file = SESSIONS_DIR / f"{session_name}.json"
    if session_file.exists():
        return f"✓ Session already saved for {session_name}. Use web_action to interact."

    pw, browser, context = await _get_browser_context(headless=True)
    page = await context.new_page()
    try:
        await page.goto(url, wait_until="domcontentloaded", timeout=30000)
        await page.wait_for_timeout(int(_jitter(500, 800) * 1000))

        # Try email first, fall back to username
        filled_user = await _smart_fill(page, "email", username)
        if not filled_user:
            filled_user = await _smart_fill(page, "username", username)

        if not filled_user:
            return f"✗ web_login: could not find username/email field on {url}"

        await page.wait_for_timeout(int(_jitter(200, 400) * 1000))
        filled_pass = await _smart_fill(page, "password", password)

        if not filled_pass:
            # Some sites show password on next screen — click continue first
            await _smart_fill(page, "submit", "")
            await page.wait_for_timeout(2000)
            filled_pass = await _smart_fill(page, "password", password)

        await page.wait_for_timeout(int(_jitter(200, 500) * 1000))
        await _smart_fill(page, "submit", "")
        await page.wait_for_timeout(3000)

        # Check if login succeeded (URL changed or dashboard element appeared)
        current_url = page.url
        title = await page.title()

        if p.get("screenshot"):
            screenshot_path = str(DATA_DIR / f"login_{session_name}.png")
            await page.screenshot(path=screenshot_path)

        path = await _save_session(context, session_name)
        # Store credentials
        accounts = _load_accounts()
        accounts[session_name] = {"url": url, "username": username, "password": password}
        _save_accounts(accounts)

        return f"✓ web_login: logged in to {url}\nSession saved: {path}\nCurrent URL: {current_url}\nTitle: {title}"

    except Exception as e:
        return f"✗ web_login error: {e}"
    finally:
        await browser.close()
        await pw.stop()


async def web_signup(p: dict) -> str:
    """
    Create a new account. Auto-handles email verification via IMAP.
    p: {url, email, password, name?, extra_fields: {label: value}, session_name?, verify_email?}
    extra_fields: any additional fields e.g. {"Username": "mybot123", "Phone": "5551234567"}
    """
    url = p.get("url", "")
    email = p.get("email", IMAP_EMAIL)
    password = p.get("password", "")
    name = p.get("name", "")
    extra_fields = p.get("extra_fields", {})
    session_name = p.get("session_name") or url.split("/")[2].replace("www.", "").replace(".", "_")
    auto_verify = p.get("verify_email", True)

    pw, browser, context = await _get_browser_context(headless=True)
    page = await context.new_page()
    try:
        await page.goto(url, wait_until="domcontentloaded", timeout=30000)
        await page.wait_for_timeout(int(_jitter(800, 1200) * 1000))

        # Fill standard fields
        if name:
            await _smart_fill(page, "name", name)
            await page.wait_for_timeout(int(_jitter() * 1000))
        await _smart_fill(page, "email", email)
        await page.wait_for_timeout(int(_jitter() * 1000))
        await _smart_fill(page, "password", password)
        await page.wait_for_timeout(int(_jitter() * 1000))

        # Fill any extra fields by placeholder/label text
        for label, value in extra_fields.items():
            try:
                el = page.locator(f'input[placeholder*="{label}" i], input[name*="{label}" i], input[id*="{label}" i]').first
                if await el.is_visible(timeout=2000):
                    await el.fill(value)
                    await page.wait_for_timeout(int(_jitter() * 1000))
            except Exception:
                pass

        await _smart_fill(page, "submit", "")
        await page.wait_for_timeout(3000)
        current_url = page.url

        result = f"✓ web_signup: form submitted on {url}\nCurrent URL: {current_url}"

        # Auto-handle email verification
        if auto_verify:
            await page.wait_for_timeout(5000)  # Wait for email to arrive
            otp_result = await extract_otp({"type": "both", "timeout_s": 30})
            if otp_result.startswith("✗"):
                result += f"\n⚠ Verification: {otp_result}"
            else:
                # otp_result contains code or link
                if "link:" in otp_result:
                    link = otp_result.split("link:")[-1].strip()
                    await page.goto(link, wait_until="domcontentloaded", timeout=20000)
                    result += f"\n✓ Clicked verification link"
                elif "code:" in otp_result:
                    code = otp_result.split("code:")[-1].strip()
                    filled = await _smart_fill(page, "otp", code)
                    if filled:
                        await _smart_fill(page, "submit", "")
                        await page.wait_for_timeout(2000)
                        result += f"\n✓ Entered OTP: {code}"
                    else:
                        result += f"\n⚠ OTP found ({code}) but could not find input field — enter manually"

        path = await _save_session(context, session_name)
        accounts = _load_accounts()
        accounts[session_name] = {"url": url, "username": email, "password": password}
        _save_accounts(accounts)
        result += f"\nSession saved: {path}"
        return result

    except Exception as e:
        return f"✗ web_signup error: {e}"
    finally:
        await browser.close()
        await pw.stop()


async def web_action(p: dict) -> str:
    """
    Execute steps on a website using a saved session.
    p: {url, session_name, steps: [...same format as automate tool...], screenshot?}
    Loads saved session so no re-login needed.
    """
    url = p.get("url", "")
    session_name = p.get("session_name", "")
    steps = p.get("steps", [])

    if not session_name:
        session_name = url.split("/")[2].replace("www.", "").replace(".", "_") if url else ""

    session_file = SESSIONS_DIR / f"{session_name}.json"
    if not session_file.exists():
        return f"✗ web_action: no saved session '{session_name}'. Run web_login first."

    pw, browser, context = await _get_browser_context(session_name=session_name, headless=True)
    page = await context.new_page()
    results = []
    try:
        await page.goto(url, wait_until="domcontentloaded", timeout=30000)
        await page.wait_for_timeout(int(_jitter(500, 800) * 1000))

        for step in steps:
            t = step.get("type", "")
            sel = step.get("selector", "")
            val = step.get("value", "")
            ms = step.get("ms", int(_jitter(300, 600) * 1000))

            try:
                if t == "fill":
                    await page.fill(sel, val)
                    results.append(f"filled {sel}")
                elif t == "type":
                    el = page.locator(sel).first
                    for char in val:
                        await el.type(char, delay=_random.randint(40, 120))
                    results.append(f"typed into {sel}")
                elif t == "click":
                    await page.click(sel)
                    await page.wait_for_timeout(int(_jitter(200, 500) * 1000))
                    results.append(f"clicked {sel}")
                elif t == "wait":
                    await page.wait_for_timeout(ms)
                    results.append(f"waited {ms}ms")
                elif t == "get_text":
                    text = await page.inner_text(sel)
                    results.append(f"text from {sel}: {text[:200]}")
                elif t == "screenshot":
                    path = val or str(DATA_DIR / f"screenshot_{session_name}_{len(results)}.png")
                    await page.screenshot(path=path)
                    results.append(f"screenshot: {path}")
                elif t == "scroll":
                    await page.evaluate(f"window.scrollBy(0, {step.get('y', 500)})")
                    results.append("scrolled")
                elif t == "url":
                    results.append(f"current url: {page.url}")
                elif t == "goto":
                    await page.goto(val, wait_until="domcontentloaded", timeout=20000)
                    results.append(f"navigated to {val}")
                elif t == "select":
                    await page.select_option(sel, val)
                    results.append(f"selected {val} in {sel}")
                elif t == "press":
                    await page.press(sel, val)
                    results.append(f"pressed {val} on {sel}")
            except Exception as e:
                results.append(f"✗ step {t} failed: {e}")

        # Re-save session after actions (refreshes cookies)
        await _save_session(context, session_name)
        return "✓ web_action complete:\n" + "\n".join(results)

    except Exception as e:
        return f"✗ web_action error: {e}"
    finally:
        await browser.close()
        await pw.stop()


async def extract_otp(p: dict) -> str:
    """
    Extract OTP code or verification link from Outlook inbox.
    p: {from_domain?, timeout_s: 30, type: "code"|"link"|"both"}
    Polls inbox until message arrives or timeout.
    """
    import imaplib, email as _email_lib, re as _re
    from email.header import decode_header

    timeout_s = int(p.get("timeout_s", 30))
    from_domain = p.get("from_domain", "")
    otp_type = p.get("type", "both")
    otp_pattern = _re.compile(r"(?<!\d)(\d{4,8})(?!\d)")
    link_pattern = _re.compile(r'https?://[^\s<>"]+(?:verify|confirm|activate|token|auth)[^\s<>"]*', _re.IGNORECASE)

    deadline = _time.time() + timeout_s
    while _time.time() < deadline:
        try:
            m = imaplib.IMAP4_SSL(IMAP_HOST, 993)
            m.login(IMAP_EMAIL, IMAP_PASS)
            m.select("inbox")
            search_q = "UNSEEN"
            if from_domain:
                search_q = f'(UNSEEN FROM "{from_domain}")'
            _, ids = m.search(None, search_q)
            msg_ids = ids[0].split()
            if msg_ids:
                # Check most recent 3 unread
                for uid in reversed(msg_ids[-3:]):
                    _, data = m.fetch(uid, "(RFC822)")
                    msg = _email_lib.message_from_bytes(data[0][1])
                    # Extract body text
                    body = ""
                    if msg.is_multipart():
                        for part in msg.walk():
                            ct = part.get_content_type()
                            if ct == "text/plain":
                                try:
                                    body += part.get_payload(decode=True).decode(errors="replace")
                                except Exception:
                                    pass
                            elif ct == "text/html" and not body:
                                try:
                                    body += part.get_payload(decode=True).decode(errors="replace")
                                except Exception:
                                    pass
                    else:
                        try:
                            body = msg.get_payload(decode=True).decode(errors="replace")
                        except Exception:
                            body = str(msg.get_payload())

                    full_text = (msg.get("Subject", "") + " " + body)[:2000]

                    # Try to find verification link
                    if otp_type in ("link", "both"):
                        link_match = link_pattern.search(full_text)
                        if link_match:
                            m.logout()
                            return f"✓ link: {link_match.group(0)}"

                    # Try to find OTP code
                    if otp_type in ("code", "both"):
                        code_match = otp_pattern.search(full_text)
                        if code_match:
                            m.logout()
                            return f"✓ code: {code_match.group(1)}"

            m.logout()
        except Exception as e:
            pass  # Retry on IMAP errors

        if _time.time() < deadline:
            _time.sleep(3)

    return f"✗ extract_otp: no verification email found within {timeout_s}s"


async def session_list(p: dict) -> str:
    """List all saved web sessions."""
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    files = list(SESSIONS_DIR.glob("*.json"))
    if not files:
        return "No saved sessions."
    accounts = _load_accounts()
    lines = []
    for f in sorted(files):
        name = f.stem
        size_kb = f.stat().st_size // 1024
        acc = accounts.get(name, {})
        lines.append(f"[{name}] {acc.get('url','')} | {acc.get('username','')} | {size_kb}KB")
    return "\n".join(lines)


async def session_delete(p: dict) -> str:
    """Delete a saved session. p: {session_name}"""
    name = p.get("session_name", "")
    if not name:
        return "✗ session_name required"
    f = SESSIONS_DIR / f"{name}.json"
    if f.exists():
        f.unlink()
        accounts = _load_accounts()
        accounts.pop(name, None)
        _save_accounts(accounts)
        return f"✓ session '{name}' deleted"
    return f"✗ session '{name}' not found"


# ── PC Control (Windows) ──────────────────────────────────────────────────

async def pc_screenshot(p: dict) -> str:
    """
    Take a screenshot of the screen or a specific window.
    p: {window: "partial title" (optional), output: "path.png" (optional)}
    """
    out = p.get("output") or os.path.join(_tempfile.gettempdir(), f"max_screen_{uuid.uuid4().hex[:8]}.png")
    window = p.get("window", "")
    script = f"""
import pyautogui, sys
pyautogui.FAILSAFE = False
window = {json.dumps(window)}
out = {json.dumps(out)}
try:
    if window:
        try:
            import win32gui, win32con
            hwnds = []
            def cb(h, _):
                if window.lower() in win32gui.GetWindowText(h).lower() and win32gui.IsWindowVisible(h):
                    hwnds.append(h)
            win32gui.EnumWindows(cb, None)
            if hwnds:
                rect = win32gui.GetWindowRect(hwnds[0])
                img = pyautogui.screenshot(region=(rect[0],rect[1],rect[2]-rect[0],rect[3]-rect[1]))
                img.save(out)
                print(out); sys.exit(0)
        except Exception:
            pass
    img = pyautogui.screenshot()
    img.save(out)
    print(out)
except Exception as e:
    print(f"error: {{e}}")
"""
    r = subprocess.run([sys.executable, "-c", script], capture_output=True, text=True,
                       timeout=15, encoding="utf-8", errors="replace")
    if r.returncode == 0 and r.stdout.strip():
        return f"✓ screenshot: {r.stdout.strip()}"
    return f"✗ pc_screenshot: {r.stderr[:300]}"


async def pc_control(p: dict) -> str:
    """
    Control the Windows PC via UI automation and keyboard/mouse.
    p: {
      action: "click"|"type"|"press"|"scroll"|"find"|"launch"|"get_windows"|"focus"|"drag",
      target: "element name or app name" or [x,y] coords,
      value: "text to type or key (ctrl+c, alt+tab, win, enter, etc.)",
      window: "window title to focus first (optional)"
    }
    Examples:
      {"action":"launch","target":"chrome"}
      {"action":"type","value":"hello world"}
      {"action":"press","value":"ctrl+c"}
      {"action":"click","target":"Submit"}
      {"action":"click","target":[960,540]}
      {"action":"get_windows"}
      {"action":"focus","window":"Chrome"}
      {"action":"find","target":"button name"}
    """
    action = p.get("action", "")
    target = p.get("target", "")
    value = p.get("value", "")
    window_hint = p.get("window", "")

    script = f"""
import sys, time, subprocess, json
import pyautogui
pyautogui.FAILSAFE = False

action = {json.dumps(action)}
target = {json.dumps(target)}
value = {json.dumps(value)}
window_hint = {json.dumps(window_hint)}

def focus_window(title):
    try:
        import win32gui, win32con
        hwnds = []
        def cb(h, _):
            if title.lower() in win32gui.GetWindowText(h).lower() and win32gui.IsWindowVisible(h):
                hwnds.append(h)
        win32gui.EnumWindows(cb, None)
        if hwnds:
            win32gui.ShowWindow(hwnds[0], win32con.SW_RESTORE)
            win32gui.SetForegroundWindow(hwnds[0])
            time.sleep(0.3)
            return win32gui.GetWindowText(hwnds[0])
    except Exception as e:
        return str(e)
    return None

try:
    if window_hint:
        focus_window(window_hint)

    if action == "get_windows":
        import win32gui
        wins = []
        def cb(h, _):
            t = win32gui.GetWindowText(h)
            if t and win32gui.IsWindowVisible(h): wins.append(t)
        win32gui.EnumWindows(cb, None)
        print("\\n".join(wins[:30]))

    elif action == "launch":
        import shutil
        app = str(target)
        found = shutil.which(app)
        if found:
            subprocess.Popen(found)
            print(f"✓ launched {{app}}")
        else:
            pyautogui.press("win"); time.sleep(0.5)
            pyautogui.write(app, interval=0.05); time.sleep(0.5)
            pyautogui.press("enter")
            print(f"✓ searched + launched {{app}}")

    elif action == "click":
        if isinstance(target, list) and len(target) == 2:
            pyautogui.click(int(target[0]), int(target[1]))
            print(f"✓ clicked at {{target[0]}},{{target[1]}}")
        else:
            import uiautomation as auto
            best, best_score = None, 0
            from difflib import SequenceMatcher
            for ctrl in auto.WalkControl(auto.GetRootControl(), maxDepth=5):
                try:
                    name = ctrl.Name or ""
                    score = SequenceMatcher(None, name.lower(), str(target).lower()).ratio()
                    if score > best_score: best_score = score; best = ctrl
                except: pass
            if best and best_score > 0.5:
                best.Click()
                print(f"✓ clicked element: {{best.Name}} (score {{best_score:.2f}})")
            else:
                print(f"✗ element not found: {{target}}")

    elif action == "type":
        pyautogui.write(str(value), interval=0.03)
        print(f"✓ typed {{len(str(value))}} chars")

    elif action == "press":
        keys = str(value).lower().replace("cmd","win").split("+")
        if len(keys) > 1: pyautogui.hotkey(*keys)
        else: pyautogui.press(keys[0])
        print(f"✓ pressed {{value}}")

    elif action == "scroll":
        clicks = int(value) if value else 3
        pyautogui.scroll(clicks)
        print(f"✓ scrolled {{clicks}}")

    elif action == "focus":
        title = str(window_hint or target)
        result = focus_window(title)
        print(f"✓ focused: {{result}}" if result else f"✗ window not found: {{title}}")

    elif action == "drag":
        coords = str(target).replace("[","").replace("]","").split(",")
        vcoords = str(value).replace("[","").replace("]","").split(",")
        if len(coords) >= 2 and len(vcoords) >= 2:
            pyautogui.dragTo(int(vcoords[0]), int(vcoords[1]), duration=0.5, mouseDownButton="left",
                             _pause=False)
            print(f"✓ dragged to {{vcoords[0]}},{{vcoords[1]}}")

    elif action == "find":
        import uiautomation as auto
        results = []
        for ctrl in auto.WalkControl(auto.GetRootControl(), maxDepth=4):
            try:
                name = ctrl.Name or ""
                if str(target).lower() in name.lower() and name:
                    results.append(f"{{ctrl.ControlTypeName}}: {{name}}")
                    if len(results) >= 10: break
            except: pass
        print("\\n".join(results) if results else f"✗ not found: {{target}}")

    else:
        print(f"✗ unknown action: {{action}}")

except Exception as e:
    print(f"✗ pc_control error: {{e}}")
"""
    r = subprocess.run([sys.executable, "-c", script], capture_output=True, text=True,
                       timeout=30, encoding="utf-8", errors="replace")
    out = (r.stdout + (f"\n[err]{r.stderr[:200]}" if r.stderr.strip() else "")).strip()
    return out or "(no output)"


async def cdp_action(p: dict) -> str:
    """
    Control Chrome via CDP (Chrome DevTools Protocol).
    Chrome must be launched with --remote-debugging-port=9222.
    Run start_chrome_cdp.bat first — it opens your real Chrome profile so all logins are inherited.

    p: {
      action: "navigate"|"click"|"fill"|"eval"|"screenshot"|"list_tabs"|"new_tab"|"get_text",
      url: "https://..." (navigate/new_tab),
      selector: "#css-selector" (click/fill/get_text),
      value: "text" or JS expression (fill/eval),
      tab_id: "..." (optional, uses first tab if omitted)
    }
    """
    import httpx, websockets, base64, time as _t

    CDP = "http://localhost:9222"
    action = p.get("action", "")

    try:
        async with httpx.AsyncClient(timeout=10) as c:
            r = await c.get(f"{CDP}/json")
            tabs = r.json()
    except Exception as e:
        return f"✗ Chrome not reachable at {CDP}. Run start_chrome_cdp.bat first. Error: {e}"

    if not tabs:
        return "✗ No Chrome tabs found."

    if action == "list_tabs":
        return "\n".join(
            f"{t.get('id','')[:8]} | {t.get('title','?')[:60]} | {t.get('url','')[:80]}"
            for t in tabs[:15]
        )

    tab_id = p.get("tab_id")
    tab = next((t for t in tabs if t.get("id","").startswith(str(tab_id))), None) if tab_id else None
    if not tab:
        tab = tabs[0]

    ws_url = tab.get("webSocketDebuggerUrl", "")
    if not ws_url:
        return "✗ No WebSocket debugger URL for this tab"

    msg_id = [1]

    async def send(ws, method, params=None):
        m = msg_id[0]; msg_id[0] += 1
        await ws.send(json.dumps({"id": m, "method": method, "params": params or {}}))
        while True:
            resp = json.loads(await asyncio.wait_for(ws.recv(), timeout=15))
            if resp.get("id") == m:
                return resp.get("result", {})

    try:
        async with websockets.connect(ws_url, max_size=10 * 1024 * 1024) as ws:
            if action == "navigate":
                url = p.get("url", "")
                await send(ws, "Page.navigate", {"url": url})
                await asyncio.sleep(2)
                return f"✓ navigated to {url}"

            elif action == "eval":
                expr = p.get("value", p.get("expression", "document.title"))
                result = await send(ws, "Runtime.evaluate", {"expression": expr, "returnByValue": True})
                val = result.get("result", {}).get("value", str(result))
                return f"✓ {json.dumps(val)[:800]}"

            elif action == "get_text":
                sel = p.get("selector", "body")
                expr = f"document.querySelector({json.dumps(sel)})?.innerText || ''"
                result = await send(ws, "Runtime.evaluate", {"expression": expr, "returnByValue": True})
                return result.get("result", {}).get("value", "")[:2000]

            elif action == "click":
                sel = p.get("selector", "")
                expr = f"document.querySelector({json.dumps(sel)})?.click()"
                await send(ws, "Runtime.evaluate", {"expression": expr})
                return f"✓ clicked {sel}"

            elif action == "fill":
                sel = p.get("selector", "")
                val = p.get("value", "")
                expr = (
                    f"var el=document.querySelector({json.dumps(sel)});"
                    f"if(el){{el.focus();el.value={json.dumps(val)};"
                    "el.dispatchEvent(new Event('input',{bubbles:true}));"
                    "el.dispatchEvent(new Event('change',{bubbles:true}));}"
                )
                await send(ws, "Runtime.evaluate", {"expression": expr})
                return f"✓ filled {sel} with {len(str(val))} chars"

            elif action == "screenshot":
                result = await send(ws, "Page.captureScreenshot", {"format": "png"})
                data = result.get("data", "")
                if data:
                    out = p.get("output") or os.path.join(_tempfile.gettempdir(), f"cdp_{int(_t.time())}.png")
                    with open(out, "wb") as f:
                        f.write(base64.b64decode(data))
                    return f"✓ screenshot: {out}"
                return "✗ no screenshot data"

            elif action == "new_tab":
                result = await send(ws, "Target.createTarget", {"url": p.get("url", "about:blank")})
                return f"✓ new tab: {result.get('targetId', '?')}"

            else:
                return f"✗ unknown CDP action: {action}"

    except Exception as e:
        return f"✗ cdp_action error: {e}"


async def runway_video(p: dict) -> str:
    """
    Generate video via RunwayML Gen-4 Turbo (working JWT).
    p: {prompt: str, image_url: str (optional for i2v), duration: 5|10, output: "path.mp4"}
    """
    import httpx, time as _time
    AUTH = f"Bearer {os.environ.get('RUNWAY_JWT', '')}"
    TEAM_ID = "56181656"
    CLIENT_ID = "1576853c-5948-4366-8e24-273b66043d6a"
    APP_VER = "da2d96b793ac79a0e7152804e4291fd07dffb252"
    BASE = "https://api.runwayml.com/v1"

    headers = {
        "accept": "application/json",
        "authorization": AUTH,
        "content-type": "application/json",
        "x-runway-client-id": CLIENT_ID,
        "x-runway-source-application": "web",
        "x-runway-source-application-version": APP_VER,
        "x-runway-workspace": TEAM_ID,
    }

    prompt = p.get("prompt", "")
    image_url = p.get("image_url", "")
    duration = int(p.get("duration", 5))
    output = p.get("output") or os.path.join(_tempfile.gettempdir(), f"runway_{uuid.uuid4().hex[:8]}.mp4")

    route = "i2v" if image_url else "t2v"
    options = {
        "route": route, "name": f"MAX-{prompt[:40]}", "text_prompt": prompt,
        "seconds": duration, "width": 1280, "height": 720,
        "exploreMode": False, "creationSource": "tool-mode", "watermark": False,
    }
    if image_url:
        options["init_image"] = image_url
        options["width"] = 720; options["height"] = 1280

    async with httpx.AsyncClient(timeout=30) as c:
        r = await c.post(f"{BASE}/tasks", headers=headers,
                         json={"taskType": "gen4_turbo", "options": options, "asTeamId": int(TEAM_ID)})
    if r.status_code not in (200, 201):
        return f"✗ runway_video: task create failed ({r.status_code}): {r.text[:300]}"

    task = r.json().get("task", r.json())
    task_id = task.get("id")
    if not task_id:
        return f"✗ runway_video: no task_id in response: {str(r.json())[:200]}"

    for _ in range(120):
        await asyncio.sleep(5)
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get(f"{BASE}/tasks/{task_id}", headers=headers, params={"asTeamId": TEAM_ID})
        task = r.json().get("task", r.json())
        status = task.get("status", "")
        if status == "SUCCEEDED":
            artifacts = task.get("artifacts", [])
            vid_url = next((a.get("url") or a.get("previewUrl") for a in artifacts if a), None)
            if not vid_url:
                return f"✗ runway_video: succeeded but no artifact URL"
            async with httpx.AsyncClient(timeout=120, follow_redirects=True) as c:
                dl = await c.get(vid_url)
            with open(output, "wb") as f:
                f.write(dl.content)
            return f"✓ runway_video: saved to {output} ({len(dl.content)//1024}KB)"
        elif status in ("FAILED", "CANCELED", "ERROR"):
            return f"✗ runway_video: {status} — {task.get('error','')}"
    return "✗ runway_video: timeout (10 min)"


async def unlimitedai(p: dict) -> str:
    """
    Free unlimited AI (no API key). Good fallback LLM when others fail.
    p: {prompt: str}
    """
    import httpx, uuid as _uuid
    from datetime import datetime, timezone as _tz
    cid = str(_uuid.uuid4()); did = str(_uuid.uuid4())
    uid = str(_uuid.uuid4()); aid = str(_uuid.uuid4())
    now = datetime.now(_tz.utc).isoformat(timespec="milliseconds").replace("+00:00","Z")
    payload = {
        "chatId": cid,
        "messages": [
            {"id":uid,"role":"user","content":p["prompt"],"parts":[{"type":"text","text":p["prompt"]}],"createdAt":now},
            {"id":aid,"role":"assistant","content":"","parts":[{"type":"text","text":""}],"createdAt":now},
        ],
        "selectedChatModel": "chat-model-reasoning",
        "selectedCharacter": None, "selectedStory": None,
        "deviceId": did, "locale": "en"
    }
    headers = {
        "accept":"*/*","content-type":"application/json",
        "referer":"https://app.unlimitedai.chat/","x-next-intl-locale":"en",
    }
    full = ""
    try:
        async with httpx.AsyncClient(timeout=60) as c:
            async with c.stream("POST","https://app.unlimitedai.chat/api/chat",
                                headers=headers, json=payload) as r:
                async for line in r.aiter_lines():
                    if not line.strip(): continue
                    try:
                        obj = json.loads(line)
                        if obj.get("type") == "delta":
                            full += obj.get("delta","")
                    except Exception:
                        pass
        return full.strip() if full else "✗ unlimitedai: empty response"
    except Exception as e:
        return f"✗ unlimitedai error: {e}"


# ── Dynamic Tool Discovery & Skill Cache ──────────────────────────────────

SKILLS_DIR = DATA_DIR / "skills"
SKILLS_DB  = str(DATA_DIR / "skills.db")

def _skills_db():
    SKILLS_DIR.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(SKILLS_DB)
    conn.row_factory = sqlite3.Row
    conn.execute("""
        CREATE TABLE IF NOT EXISTS skills (
            name        TEXT PRIMARY KEY,
            category    TEXT DEFAULT '',
            description TEXT DEFAULT '',
            auth_level  INTEGER DEFAULT 3,
            endpoint    TEXT DEFAULT '',
            code        TEXT DEFAULT '',
            usage_count INTEGER DEFAULT 0,
            last_used   TEXT,
            verified    INTEGER DEFAULT 0,
            created_at  TEXT NOT NULL
        )
    """)
    conn.commit()
    return conn


async def skill_list(p: dict) -> str:
    """
    List all cached skills/tools.
    p: {category: "image|video|llm|tts|search|ocr|..." (optional), query: "keyword"}
    """
    cat = p.get("category", "")
    query = p.get("query", "").lower()
    with _skills_db() as conn:
        if cat:
            rows = conn.execute(
                "SELECT * FROM skills WHERE category=? ORDER BY usage_count DESC",
                (cat,)
            ).fetchall()
        elif query:
            rows = conn.execute(
                "SELECT * FROM skills WHERE lower(name)||lower(description)||lower(category) LIKE ? ORDER BY usage_count DESC",
                (f"%{query}%",)
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT * FROM skills ORDER BY usage_count DESC LIMIT 50"
            ).fetchall()
    if not rows:
        return "No skills cached yet. Run tool_discover to find and save tools."
    lines = []
    for r in rows:
        auth = ["no-auth", "free-key", "signup", "paid"][min(r["auth_level"], 3)]
        lines.append(
            f"[{r['name']}] {r['category']} | {auth} | used {r['usage_count']}x | {r['description'][:80]}"
        )
    return "\n".join(lines)


async def skill_save(p: dict) -> str:
    """
    Save a discovered/verified tool as a reusable skill.
    p: {name, category, description, auth_level: 0-3, endpoint, code, verified: true}
    auth_level: 0=no-auth, 1=free-api-key, 2=signup-required, 3=paid
    """
    name = p.get("name", "").strip().replace(" ", "_").lower()
    if not name:
        return "✗ name required"
    now = datetime.now(timezone.utc).isoformat()
    with _skills_db() as conn:
        conn.execute("""
            INSERT OR REPLACE INTO skills
            (name,category,description,auth_level,endpoint,code,verified,created_at)
            VALUES (?,?,?,?,?,?,?,?)
        """, (
            name,
            p.get("category", "general"),
            p.get("description", ""),
            int(p.get("auth_level", 2)),
            p.get("endpoint", ""),
            p.get("code", ""),
            1 if p.get("verified") else 0,
            now,
        ))
    # Also write code to skills dir for direct execution
    if p.get("code"):
        skill_file = SKILLS_DIR / f"{name}.py"
        skill_file.write_text(p["code"], encoding="utf-8")
    return f"✓ skill saved: {name}"


async def skill_run(p: dict) -> str:
    """
    Run a cached skill by name.
    p: {name: "skill_name", params: {...any params the skill needs}}
    """
    name = p.get("name", "").replace(" ", "_").lower()
    params = p.get("params", p.get("input", {}))
    with _skills_db() as conn:
        row = conn.execute("SELECT * FROM skills WHERE name=?", (name,)).fetchone()
        if not row:
            return f"✗ skill '{name}' not found. Run skill_list to see available skills."
        # Update usage count
        conn.execute(
            "UPDATE skills SET usage_count=usage_count+1, last_used=? WHERE name=?",
            (datetime.now(timezone.utc).isoformat(), name)
        )
    code = row["code"]
    if not code:
        skill_file = SKILLS_DIR / f"{name}.py"
        if skill_file.exists():
            code = skill_file.read_text(encoding="utf-8")
    if not code:
        return f"✗ skill '{name}' has no executable code. Endpoint: {row['endpoint']}"

    # Inject params and run
    runner = f"""
import asyncio, json, sys, os
sys.path.insert(0, {json.dumps(str(Path(__file__).parent))})
params = {json.dumps(params)}

{code}

# Auto-detect entry point
import inspect
_fns = [(k,v) for k,v in list(locals().items()) + list(globals().items())
        if callable(v) and not k.startswith('_') and k not in ('asyncio','json','sys','os','inspect')]
_fn = next((v for k,v in _fns if 'run' in k or 'execute' in k or 'main' in k or 'call' in k), None)
if _fn is None and _fns:
    _fn = _fns[-1][1]
if _fn:
    result = _fn(params) if not asyncio.iscoroutinefunction(_fn) else asyncio.run(_fn(params))
    print(json.dumps(result) if not isinstance(result, str) else result)
else:
    print("(no entry point found)")
"""
    r = subprocess.run(
        [sys.executable, "-c", runner],
        capture_output=True, text=True, timeout=60,
        encoding="utf-8", errors="replace"
    )
    out = (r.stdout + (f"\n[err]{r.stderr[:300]}" if r.stderr.strip() else "")).strip()
    return out or "(no output)"


async def tool_discover(p: dict) -> str:
    """
    Discover the best free tool for a task. Searches in parallel across:
    - public-apis GitHub repo (no-auth APIs)
    - APIs.guru directory
    - Web search for free APIs
    - Known built-in no-auth endpoints

    Scores by: 0=no-auth (best) → 1=free-key → 2=signup → 3=paid (avoid)

    p: {
      task: "what you want to do (e.g. 'generate image', 'translate text', 'speech to text')",
      category: "image|video|llm|tts|search|ocr|translate|weather|..." (optional hint),
      save: true (auto-save best result as skill)
    }
    Returns: top 3 candidates with endpoints, auth requirements, and ready-to-paste code.
    """
    import httpx

    task = p.get("task", p.get("query", ""))
    category = p.get("category", "")
    auto_save = p.get("save", False)

    if not task:
        return "✗ task required. Example: tool_discover {\"task\":\"generate image free\"}"

    # ── Built-in no-auth endpoints (auth_level=0) ────────────────────────
    KNOWN_FREE = [
        {"name":"open_meteo_weather","category":"weather","auth_level":0,
         "endpoint":"https://api.open-meteo.com/v1/forecast",
         "description":"Real-time weather forecast, no key. Params: latitude, longitude",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient() as c:
        r = await c.get("https://api.open-meteo.com/v1/forecast",
            params={"latitude":params.get("lat",48.85),"longitude":params.get("lon",2.35),"current_weather":True})
    return r.json()["current_weather"]'''},
        {"name":"restcountries","category":"geo","auth_level":0,
         "endpoint":"https://restcountries.com/v3.1/name/{name}",
         "description":"Country info by name, no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient() as c:
        r = await c.get(f"https://restcountries.com/v3.1/name/{params['name']}")
    return r.json()[0]'''},
        {"name":"ipify_geoip","category":"network","auth_level":0,
         "endpoint":"https://api.ipify.org?format=json",
         "description":"Get public IP, no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient() as c:
        r = await c.get("https://api.ipify.org?format=json")
    return r.json()'''},
        {"name":"exchangerate_api","category":"finance","auth_level":0,
         "endpoint":"https://open.er-api.com/v6/latest/{base}",
         "description":"Live currency exchange rates, no key",
         "code":'''async def run(params):
    import httpx
    base = params.get("base","USD")
    async with httpx.AsyncClient() as c:
        r = await c.get(f"https://open.er-api.com/v6/latest/{base}")
    return r.json()["rates"]'''},
        {"name":"duckduckgo_search","category":"search","auth_level":0,
         "endpoint":"https://html.duckduckgo.com/html/",
         "description":"Web search, no key, no signup",
         "code":'''async def run(params):
    import httpx, re
    async with httpx.AsyncClient(follow_redirects=True) as c:
        r = await c.get("https://html.duckduckgo.com/html/",
            params={"q":params["query"]},
            headers={"User-Agent":"Mozilla/5.0"})
    results = re.findall(r\'class="result__a"[^>]*>(.*?)</a>.*?class="result__snippet"[^>]*>(.*?)</a>\',r.text,re.DOTALL)
    return [{"title":re.sub("<[^>]+>","",t).strip(),"snippet":re.sub("<[^>]+>","",s).strip()} for t,s in results[:5]]'''},
        {"name":"wikipedia_summary","category":"knowledge","auth_level":0,
         "endpoint":"https://en.wikipedia.org/api/rest_v1/page/summary/{title}",
         "description":"Wikipedia page summary, no key",
         "code":'''async def run(params):
    import httpx
    title = params.get("title","").replace(" ","_")
    async with httpx.AsyncClient() as c:
        r = await c.get(f"https://en.wikipedia.org/api/rest_v1/page/summary/{title}")
    return r.json().get("extract","not found")'''},
        {"name":"catfact","category":"fun","auth_level":0,
         "endpoint":"https://catfact.ninja/fact","description":"Random cat fact, no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient() as c:
        r = await c.get("https://catfact.ninja/fact")
    return r.json()["fact"]'''},
        {"name":"qrcode_gen","category":"image","auth_level":0,
         "endpoint":"https://api.qrserver.com/v1/create-qr-code/",
         "description":"Generate QR code PNG, no key. Param: data (text/URL)",
         "code":'''async def run(params):
    url = f"https://api.qrserver.com/v1/create-qr-code/?size=300x300&data={params.get(\'data\',\'hello\')}"
    return url'''},
        {"name":"placeholder_image","category":"image","auth_level":0,
         "endpoint":"https://picsum.photos/{width}/{height}",
         "description":"Random placeholder image, no key",
         "code":'''async def run(params):
    w,h = params.get("width",800),params.get("height",600)
    return f"https://picsum.photos/{w}/{h}"'''},
        {"name":"translate_mymemory","category":"translate","auth_level":0,
         "endpoint":"https://api.mymemory.translated.net/get",
         "description":"Text translation, no key (1000 chars/day free)",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient() as c:
        r = await c.get("https://api.mymemory.translated.net/get",
            params={"q":params["text"],"langpair":f"{params.get(\'from\',\'en\')}|{params.get(\'to\',\'fr\')}"})
    return r.json()["responseData"]["translatedText"]'''},
        {"name":"pdf_to_text_api","category":"ocr","auth_level":0,
         "endpoint":"https://api.extractpdf.com/",
         "description":"Extract text from PDF URL, no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient(timeout=30) as c:
        r = await c.post("https://api.extractpdf.com/extract",
            json={"url":params["url"]})
    return r.json().get("text","")'''},
        {"name":"html_to_md","category":"convert","auth_level":0,
         "endpoint":"https://md.dhr.wtf/",
         "description":"Convert any URL to Markdown, no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient(timeout=20) as c:
        r = await c.get(f"https://md.dhr.wtf/?url={params[\'url\']}")
    return r.text[:3000]'''},
        {"name":"hacker_news_top","category":"news","auth_level":0,
         "endpoint":"https://hacker-news.firebaseio.com/v0/topstories.json",
         "description":"Hacker News top stories, no key",
         "code":'''async def run(params):
    import httpx
    n = params.get("count", 5)
    async with httpx.AsyncClient() as c:
        ids = (await c.get("https://hacker-news.firebaseio.com/v0/topstories.json")).json()[:n]
        stories = []
        for i in ids:
            s = (await c.get(f"https://hacker-news.firebaseio.com/v0/item/{i}.json")).json()
            stories.append({"title":s.get("title"),"url":s.get("url"),"score":s.get("score")})
    return stories'''},
        {"name":"vondy_video_gen","category":"video","auth_level":0,
         "endpoint":"https://api.apivondy.com/api/open/video/generate",
         "description":"AI video generation, no signup, no key (Vondy)",
         "code":'''async def run(params):
    import httpx, uuid, random, time
    ua = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/147.0.0.0 Safari/537.36"
    headers = {"Content-Type":"application/json","User-Agent":ua,
               "x-vondy-tier":"lite","x-device-id":f"{uuid.uuid4().hex[:16]}-{random.getrandbits(32):x}",
               "Referer":"https://vondy.com/","Origin":"https://vondy.com/"}
    async with httpx.AsyncClient(timeout=45) as c:
        r = await c.post("https://api.apivondy.com/api/open/video/generate",
            headers=headers,json={"prompt":params["prompt"],"duration":"5","aspectRatio":"16:9","resolution":"720p","cameraFixed":False})
    data = r.json()
    return data.get("url") or data.get("resultUrl") or data.get("videoUrl") or str(data)'''},
        {"name":"freegen_image","category":"image","auth_level":0,
         "endpoint":"https://image-generator.freegen.app/",
         "description":"Free image generation via freegen.app, no key",
         "code":'''async def run(params):
    import httpx, json, hashlib, base64, time, websockets, asyncio
    sign = (await (httpx.AsyncClient()).post("https://prompt-signer.freegen.app/",json={"prompt":params["prompt"]})).json()
    job = (await (httpx.AsyncClient()).post("https://image-generator.freegen.app/",
        json={"prompt":params["prompt"],"ts":sign["ts"],"sig":sign["sig"],"ratio_id":params.get("ratio","16:9")})).json()
    jid = job["job_id"]
    ts = int(time.time()); msg = f"{jid}{ts}"; auth = base64.b64encode(hashlib.sha256(msg.encode()).hexdigest().encode()).decode()[:20]+":"+str(ts)
    async with websockets.connect("wss://websocket-bridge.freegen.app/ws") as ws:
        await ws.send(json.dumps({"type":"subscribe","job_id":jid,"auth":auth}))
        while True:
            m = json.loads(await asyncio.wait_for(ws.recv(),timeout=60))
            if m.get("type") == "result": return m.get("image_data","")[:200]+"..."
'''},
        {"name":"unlimitedai_llm","category":"llm","auth_level":0,
         "endpoint":"https://app.unlimitedai.chat/api/chat",
         "description":"Free unlimited LLM (reasoning model), no key, no signup",
         "code":'''async def run(params):
    import httpx, uuid, json
    from datetime import datetime, timezone
    cid,did,uid,aid = str(uuid.uuid4()),str(uuid.uuid4()),str(uuid.uuid4()),str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat(timespec="milliseconds").replace("+00:00","Z")
    payload = {"chatId":cid,"messages":[
        {"id":uid,"role":"user","content":params["prompt"],"parts":[{"type":"text","text":params["prompt"]}],"createdAt":now},
        {"id":aid,"role":"assistant","content":"","parts":[{"type":"text","text":""}],"createdAt":now}],
        "selectedChatModel":"chat-model-reasoning","selectedCharacter":None,"selectedStory":None,"deviceId":did,"locale":"en"}
    full=""
    async with httpx.AsyncClient(timeout=60) as c:
        async with c.stream("POST","https://app.unlimitedai.chat/api/chat",
            headers={"content-type":"application/json","referer":"https://app.unlimitedai.chat/"},json=payload) as r:
            async for line in r.aiter_lines():
                try:
                    obj=json.loads(line)
                    if obj.get("type")=="delta": full+=obj.get("delta","")
                except: pass
    return full.strip()'''},
        {"name":"glm_llm","category":"llm","auth_level":0,
         "endpoint":"https://opus4-6.online/api/chat",
         "description":"Free GLM LLM, no key, no signup",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient(timeout=30) as c:
        r = await c.post("https://opus4-6.online/api/chat",
            json={"messages":[{"role":"user","content":params["prompt"]}]},
            headers={"content-type":"application/json"})
    return r.json().get("choices",[{}])[0].get("message",{}).get("content","")'''},
        {"name":"edge_tts","category":"tts","auth_level":0,
         "endpoint":"edge-tts Python library",
         "description":"Microsoft Edge TTS, no key, 100+ voices",
         "code":'''async def run(params):
    import edge_tts, tempfile, os
    text = params.get("text","Hello")
    voice = params.get("voice","en-US-JennyNeural")
    out = params.get("output", tempfile.mktemp(suffix=".mp3"))
    comm = edge_tts.Communicate(text, voice)
    await comm.save(out)
    return out'''},
        {"name":"searxng_search","category":"search","auth_level":0,
         "endpoint":"http://localhost:4040/search",
         "description":"Self-hosted SearXNG (20+ engines), no key",
         "code":'''async def run(params):
    import httpx
    async with httpx.AsyncClient(timeout=15) as c:
        r = await c.get("http://localhost:4040/search",
            params={"q":params["query"],"format":"json","pageno":1})
    return [{"title":x.get("title"),"url":x.get("url"),"snippet":x.get("content","")} for x in r.json().get("results",[])[:8]]'''},
        {"name":"trafilatura_scrape","category":"scrape","auth_level":0,
         "endpoint":"trafilatura Python library",
         "description":"Extract clean text from any URL, no key",
         "code":'''async def run(params):
    import httpx, subprocess, sys
    r = subprocess.run([sys.executable,"-c",
        f"import trafilatura; print(trafilatura.extract(trafilatura.fetch_url({repr(params[\'url\'])})) or \'failed\')"],
        capture_output=True,text=True,timeout=20)
    return r.stdout.strip()[:3000]'''},
    ]

    # ── Score and filter built-in known tools by task relevance ─────────────
    task_lower = task.lower()
    category_lower = category.lower()

    def _relevance(tool: dict) -> int:
        score = 0
        haystack = (tool["name"] + tool["description"] + tool["category"]).lower()
        for word in task_lower.split() + category_lower.split():
            if word and len(word) > 2 and word in haystack:
                score += 1
        score += (3 - tool["auth_level"]) * 2  # heavily favor no-auth
        return score

    local_hits = sorted(
        [t for t in KNOWN_FREE if _relevance(t) > 0],
        key=_relevance, reverse=True
    )[:3]

    # ── Parallel web search for unknown tasks ────────────────────────────────
    web_results = []
    if not local_hits or len(local_hits) < 2:
        queries = [
            f"free API no signup no auth {task} site:rapidapi.com OR site:github.com",
            f"free no key API {task} json endpoint",
            f"public API no authentication {task}",
        ]
        async def _search_one(q):
            try:
                return await search({"query": q, "num_results": 5})
            except Exception:
                return ""

        results = await asyncio.gather(*[_search_one(q) for q in queries])
        web_results = [r for r in results if r]

    # ── LLM synthesis: pick best option ────────────────────────────────────
    synthesis_prompt = f"""You are a tool discovery engine. Task: "{task}"

Known no-auth tools that may match:
{json.dumps([{"name":t["name"],"description":t["description"],"endpoint":t["endpoint"]} for t in local_hits], indent=2)}

Web search findings:
{chr(10).join(web_results[:2])[:1500]}

Your job:
1. Pick the TOP 3 best free tools for this task. Prefer: no-auth > free-key > signup.
2. For each, output READY-TO-PASTE Python async code using httpx.
3. Include actual working endpoint URLs.

Respond as JSON:
{{
  "tools": [
    {{
      "name": "snake_case_name",
      "category": "image|video|llm|tts|search|ocr|translate|weather|other",
      "auth_level": 0,
      "description": "what it does",
      "endpoint": "https://...",
      "code": "async def run(params):\\n    ...",
      "why": "why this is the best option"
    }}
  ]
}}"""

    llm_raw = await llm({"prompt": synthesis_prompt, "max_tokens": 1500})

    # Parse LLM response
    discovered = []
    try:
        m = re.search(r'\{[\s\S]*\}', llm_raw)
        if m:
            parsed = json.loads(m.group())
            discovered = parsed.get("tools", [])
    except Exception:
        pass

    # Merge local hits (they take priority — they're verified)
    all_tools = local_hits + [
        t for t in discovered
        if not any(t.get("name","") == l["name"] for l in local_hits)
    ]
    all_tools = all_tools[:3]

    if not all_tools:
        return f"✗ No tools found for '{task}'. Try a more specific task."

    # Auto-save best tool as skill
    if auto_save and all_tools:
        best = all_tools[0]
        await skill_save({
            "name": best.get("name","discovered_tool"),
            "category": best.get("category","general"),
            "description": best.get("description",""),
            "auth_level": best.get("auth_level", 2),
            "endpoint": best.get("endpoint",""),
            "code": best.get("code",""),
            "verified": False,
        })

    # Format output
    lines = [f"TOOL DISCOVERY: '{task}'\n{'='*50}"]
    for i, t in enumerate(all_tools, 1):
        auth_label = ["NO-AUTH (best)", "FREE API KEY", "SIGNUP REQUIRED", "PAID"][min(t.get("auth_level",2), 3)]
        lines.append(f"\n#{i} [{auth_label}] {t.get('name','?')}")
        lines.append(f"   Category : {t.get('category','?')}")
        lines.append(f"   Endpoint : {t.get('endpoint','?')}")
        lines.append(f"   What     : {t.get('description','?')}")
        if t.get("why"):
            lines.append(f"   Why best : {t.get('why')}")
        if t.get("code"):
            lines.append(f"   Code:\n```python\n{t['code'][:600]}\n```")
    lines.append(f"\n→ Save best: python tools.py skill_save '{{\"name\":\"{all_tools[0].get('name','tool')}\",\"category\":\"{all_tools[0].get('category','general')}\",\"auth_level\":{all_tools[0].get('auth_level',0)},\"endpoint\":\"{all_tools[0].get('endpoint','')}\",\"code\":\"...\",\"description\":\"{all_tools[0].get('description','')}\"}}'")
    lines.append(f"→ Run best : python tools.py skill_run '{{\"name\":\"{all_tools[0].get('name','tool')}\",\"params\":{{...}}}}'")
    return "\n".join(lines)


# ── Registry ───────────────────────────────────────────────────────────────
TOOLS = {
    # Task tracking
    "task_create":        task_create,
    "task_update":        task_update,
    "task_list":          task_list,
    # Micro-tasking (RSTD)
    "plan_task":          plan_task,
    "subtask_create":     subtask_create,
    "subtask_update":     subtask_update,
    "subtask_list":       subtask_list,
    "validate_result":    validate_result,
    "escalate_task":      escalate_task,
    # Communication
    "send_email":         send_email,
    "read_emails":        read_emails,
    "send_telegram":      send_telegram,
    "send_telegram_photo": send_telegram_photo,
    "send_telegram_document": send_telegram_document,
    "send_whatsapp":      send_whatsapp,
    "host_file":          host_file,
    # Media
    "generate_image":     generate_image,
    "generate_video":     generate_video,
    "generate_video_free": generate_video_free,
    "generate_powerpoint": generate_powerpoint,
    # Web
    "search":             search,
    "scrape":             scrape,
    "browse":             browse,
    "automate":           automate,
    # LLM
    "llm":                llm,
    # Web auth & session management
    "web_login":          web_login,
    "web_signup":         web_signup,
    "web_action":         web_action,
    "extract_otp":        extract_otp,
    "session_list":       session_list,
    "session_delete":     session_delete,
    # Dynamic tool discovery & skill cache
    "tool_discover":      tool_discover,
    "skill_save":         skill_save,
    "skill_list":         skill_list,
    "skill_run":          skill_run,
    # PC control (Windows)
    "pc_screenshot":      pc_screenshot,
    "pc_control":         pc_control,
    # Chrome CDP bridge
    "cdp_action":         cdp_action,
    # Video generation
    "vondy_video":        vondy_video,
    "runway_video":       runway_video,
    # Free LLM backup
    "unlimitedai":        unlimitedai,
}


async def main():
    if len(sys.argv) < 2:
        print("Tools:", ", ".join(TOOLS))
        return
    name = sys.argv[1]
    params = json.loads(sys.argv[2]) if len(sys.argv) > 2 else {}
    if name not in TOOLS:
        print(f"Unknown: {name}. Available: {', '.join(TOOLS)}")
        return
    print(await TOOLS[name](params))


if __name__ == "__main__":
    asyncio.run(main())
