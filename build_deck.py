# -*- coding: utf-8 -*-
import sys, subprocess, os, urllib.request, io
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2  = RGBColor(0x12, 0x2C, 0x52)
GOLD   = RGBColor(0xE8, 0xB5, 0x4B)
WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
GREY   = RGBColor(0xB8, 0xC4, 0xD4)
GREEN  = RGBColor(0x3D, 0xC9, 0x7A)
CARD   = RGBColor(0x16, 0x33, 0x5C)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# ---- download images ----
IMGS = {
    "cover":  "https://max.vdo-x.art/img/img_1781281608936.jpg",
    "growth": "https://max.vdo-x.art/img/img_1781281608225.jpg",
    "pie":    "https://max.vdo-x.art/img/img_1781281608005.jpg",
    "invest": "https://max.vdo-x.art/img/img_1781281608888.jpg",
}
os.makedirs(r"C:\tmp\deck", exist_ok=True)
local = {}
for k, u in IMGS.items():
    fp = rf"C:\tmp\deck\{k}.jpg"
    try:
        req = urllib.request.Request(u, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=40) as r:
            data = r.read()
        if len(data) > 1500:
            with open(fp, "wb") as f:
                f.write(data)
            local[k] = fp
    except Exception as e:
        print("img fail", k, e)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=6.0):
    """runs: list of paragraphs; each paragraph = list of (txt,size,color,bold)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Segoe UI"
    return tb

def header(slide, kicker, title):
    rect(slide, 0, 0, Inches(13.333), Inches(0.18), GOLD)
    rect(slide, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.95), GOLD)
    text(slide, Inches(0.85), Inches(0.45), Inches(11.5), Inches(1.2),
         [[(kicker.upper(), 13, GOLD, True)],
          [(title, 32, WHITE, True)]], space=2.0)

def bullets(slide, x, y, w, items, size=16, gap=10.0):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            runs.append([("▸  ", size, GOLD, True), (head, size, WHITE, True),
                         ("  —  " + body, size, GREY, False)])
        else:
            runs.append([("▸  ", size, GOLD, True), (it, size, WHITE, False)])
    text(slide, x, y, w, Inches(5), runs, space=gap)

def card(slide, x, y, w, h, title, body, tcolor=GOLD):
    rect(slide, x, y, w, h, CARD)
    rect(slide, x, y, w, Inches(0.08), tcolor)
    text(slide, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), h - Inches(0.4),
         [[(title, 17, tcolor, True)], [(body, 13, GREY, False)]], space=5.0)

def pic_cover(slide, fp, x, y, w, h):
    """add picture cropped to fill the box (cover)."""
    from PIL import Image
    try:
        iw, ih = Image.open(fp).size
    except Exception:
        slide.shapes.add_picture(fp, x, y, w, h); return
    box_ratio = w / h
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(fp, x, y, width=w, height=h)
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic

# ============ SLIDE 1 — COVER ============
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
if "cover" in local:
    pic_cover(s, local["cover"], 0, 0, EMU_W, EMU_H)
# dark overlay for legibility
ov = rect(s, 0, 0, EMU_W, EMU_H, NAVY)
ov.fill.fore_color.rgb = NAVY
ov.fill.transparency = 0  # set via xml below
# transparency hack
from pptx.oxml.ns import qn
sp = ov.fill._xPr.find(qn('a:solidFill'))
srgb = sp.find(qn('a:srgbClr'))
alpha = srgb.makeelement(qn('a:alpha'), {'val': '42000'})
srgb.append(alpha)
rect(s, Inches(0.7), Inches(2.5), Inches(0.18), Inches(2.0), GOLD)
text(s, Inches(1.05), Inches(2.3), Inches(11), Inches(3),
     [[("THE STOCK MARKET", 16, GOLD, True)],
      [("Investing the Smart Way", 54, WHITE, True)],
      [("A clear, practical guide to building long-term wealth", 22, GREY, False)]],
     space=10.0)
text(s, Inches(1.05), Inches(6.6), Inches(11), Inches(0.6),
     [[("Prepared by MAX  ·  Personal Finance Briefing", 13, GREY, False)]])

# ============ SLIDE 2 — WHAT IS IT ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Foundations", "What Is the Stock Market?")
text(s, Inches(0.85), Inches(1.9), Inches(7.0), Inches(4.5),
     [[("A stock market is a network of exchanges where shares of public "
        "companies are bought and sold.", 17, WHITE, False)],
      [("When you buy a share, you own a small slice of a real business "
        "and a claim on its future profits.", 15, GREY, False)]], space=12.0)
bullets(s, Inches(0.85), Inches(3.7), Inches(7.0), [
    ("Exchanges", "NYSE & Nasdaq match buyers and sellers"),
    ("Indexes", "S&P 500, Dow & Nasdaq track the market's health"),
    ("Price", "set continuously by supply, demand & expectations"),
    ("Returns", "come from price growth plus dividends"),
], size=15, gap=9.0)
card(s, Inches(8.4), Inches(2.0), Inches(4.1), Inches(2.0),
     "~10% / year", "Historical average annual return of the S&P 500 before inflation, over the long run.", GREEN)
card(s, Inches(8.4), Inches(4.3), Inches(4.1), Inches(2.0),
     "Ownership, not lottery", "Investing is buying productive businesses — not gambling on tickers.", GOLD)

# ============ SLIDE 3 — WHY INVEST ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Motivation", "Why Invest At All?")
if "growth" in local:
    pic_cover(s, local["growth"], Inches(7.9), Inches(1.9), Inches(4.8), Inches(4.6))
bullets(s, Inches(0.85), Inches(2.0), Inches(6.7), [
    ("Beat inflation", "cash loses ~2-3% of value every year"),
    ("Compound growth", "your returns earn their own returns"),
    ("Build wealth", "small, regular amounts become large sums"),
    ("Reach goals", "retirement, a home, freedom of choice"),
    ("Stay ahead", "money parked in savings barely keeps up"),
], size=17, gap=14.0)

# ============ SLIDE 4 — BUILDING BLOCKS ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "The Toolkit", "Core Building Blocks")
defs = [
    ("Stocks", "Ownership in a single company. Higher risk, higher potential reward."),
    ("Bonds", "A loan to a government or company. Steadier, lower returns."),
    ("Index Funds", "Own an entire market in one low-cost fund. Best for beginners."),
    ("ETFs", "Funds that trade like stocks. Flexible, cheap, diversified."),
    ("Mutual Funds", "Professionally managed baskets of assets."),
    ("Dividends", "Cash a company pays you just for holding its shares."),
]
xs = [Inches(0.85), Inches(5.1), Inches(9.35)]
ys = [Inches(2.1), Inches(4.45)]
for i, (t, b) in enumerate(defs):
    card(s, xs[i % 3], ys[i // 3], Inches(3.8), Inches(2.05), t, b)

# ============ SLIDE 5 — COMPOUNDING ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "The Secret Weapon", "The Power of Compounding")
text(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(0.9),
     [[("Invest ", 18, WHITE, False), ("$300 / month", 18, GOLD, True),
       (" at a ", 18, WHITE, False), ("8% average annual return", 18, GOLD, True),
       (":", 18, WHITE, False)]])
cols = [("After 10 yrs", "$54,900", "you put in $36,000"),
        ("After 20 yrs", "$176,700", "you put in $72,000"),
        ("After 30 yrs", "$447,100", "you put in $108,000"),
        ("After 40 yrs", "$1,048,000", "you put in $144,000")]
for i, (h, v, sub) in enumerate(cols):
    x = Inches(0.85 + i * 3.05)
    rect(s, x, Inches(3.0), Inches(2.8), Inches(2.7), CARD)
    rect(s, x, Inches(3.0), Inches(2.8), Inches(0.08), GREEN)
    text(s, x, Inches(3.2), Inches(2.8), Inches(2.4),
         [[(h, 14, GREY, True)], [(v, 30, GREEN, True)], [(sub, 12, GREY, False)]],
         align=PP_ALIGN.CENTER, space=8.0)
text(s, Inches(0.85), Inches(6.1), Inches(11.6), Inches(0.8),
     [[("The earlier you start, the more the math works for you. Time beats timing.",
        16, GOLD, True)]], align=PP_ALIGN.CENTER)

# ============ SLIDE 6 — RISK & DIVERSIFICATION ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Protect Your Downside", "Risk & Diversification")
if "pie" in local:
    pic_cover(s, local["pie"], Inches(8.1), Inches(2.0), Inches(4.6), Inches(4.5))
bullets(s, Inches(0.85), Inches(2.0), Inches(6.9), [
    ("Don't bet on one horse", "spread money across many companies & sectors"),
    ("Mix asset types", "stocks, bonds & cash balance growth and safety"),
    ("Go global", "different countries rise and fall at different times"),
    ("Index funds = instant diversification", "hundreds of stocks in one buy"),
    ("Match risk to time", "the longer your horizon, the more risk you can hold"),
], size=15, gap=11.0)

# ============ SLIDE 7 — ASSET ALLOCATION ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Building Your Mix", "Asset Allocation by Life Stage")
rows = [
    ("20s - 30s  ·  Aggressive", "90% Stocks", "10% Bonds", "Decades to recover — maximize growth."),
    ("40s - 50s  ·  Balanced", "70% Stocks", "30% Bonds", "Still growing, but adding stability."),
    ("60s+  ·  Conservative", "50% Stocks", "50% Bonds", "Protect capital, generate income."),
]
y = Inches(2.1)
for label, a, b, note in rows:
    rect(s, Inches(0.85), y, Inches(11.6), Inches(1.35), CARD)
    text(s, Inches(1.1), y + Inches(0.18), Inches(4.0), Inches(1.0),
         [[(label, 18, WHITE, True)], [(note, 12.5, GREY, False)]], space=4.0)
    text(s, Inches(7.6), y + Inches(0.3), Inches(2.3), Inches(0.8),
         [[(a, 22, GREEN, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(10.0), y + Inches(0.3), Inches(2.3), Inches(0.8),
         [[(b, 22, GOLD, True)]], align=PP_ALIGN.CENTER)
    y = Inches(y.inches + 1.6)
text(s, Inches(0.85), Inches(6.95), Inches(11.6), Inches(0.4),
     [[("Rule of thumb: subtract your age from 110 for your % in stocks.", 13, GREY, False)]])

# ============ SLIDE 8 — STRATEGIES ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "How Winners Play", "Smart Investing Strategies")
strat = [
    ("Dollar-Cost Averaging", "Invest a fixed amount on a schedule. Removes emotion and timing risk."),
    ("Buy & Hold", "Stay invested for years. Time in the market beats timing the market."),
    ("Index Investing", "Own the whole market cheaply. Beats most active funds over time."),
    ("Reinvest Dividends", "Turn payouts back into shares to supercharge compounding."),
    ("Stay the Course", "Don't panic-sell in downturns — they are when wealth is built."),
    ("Keep Costs Low", "Fees quietly eat returns. Favor funds under 0.20% expense ratio."),
]
xs = [Inches(0.85), Inches(6.7)]
ys = [Inches(2.05), Inches(3.6), Inches(5.15)]
for i, (t, b) in enumerate(strat):
    card(s, xs[i % 2], ys[i // 2], Inches(5.75), Inches(1.4), t, b, GOLD)

# ============ SLIDE 9 — MISTAKES ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Stay Out of Trouble", "Common Mistakes to Avoid")
RED = RGBColor(0xE5, 0x6B, 0x6B)
mistakes = [
    "Trying to time the market — almost no one does it consistently",
    "Panic-selling when prices drop and locking in losses",
    "Chasing hype, meme stocks, and 'get rich quick' tips",
    "Putting everything into one stock or one sector",
    "Paying high fees and ignoring their long-term drag",
    "Investing money you'll need within 1-2 years",
    "Checking your portfolio daily and trading on emotion",
]
runs = [[("✕  ", 16, RED, True), (m, 16, WHITE, False)] for m in mistakes]
text(s, Inches(0.85), Inches(2.0), Inches(11.6), Inches(5), runs, space=12.0)

# ============ SLIDE 10 — GET STARTED ============
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Your Action Plan", "How to Get Started")
if "invest" in local:
    pic_cover(s, local["invest"], Inches(8.1), Inches(1.95), Inches(4.6), Inches(4.6))
steps = [
    ("1. Build a buffer", "Save 3-6 months of expenses before investing"),
    ("2. Kill high-interest debt", "Pay off credit cards first — guaranteed return"),
    ("3. Open an account", "A brokerage or retirement account (401k / IRA)"),
    ("4. Automate it", "Set a fixed monthly auto-invest and forget it"),
    ("5. Buy a broad index fund", "Low-cost S&P 500 or total-market ETF"),
    ("6. Stay consistent", "Keep investing through ups and downs"),
]
runs = []
for h, b in steps:
    runs.append([(h, 16, GOLD, True), ("  —  " + b, 14, GREY, False)])
text(s, Inches(0.85), Inches(2.05), Inches(6.9), Inches(5), runs, space=13.0)

# ============ SLIDE 11 — TAKEAWAYS ============
s = prs.slides.add_slide(BLANK); bg(s, NAVY2)
rect(s, 0, 0, EMU_W, Inches(0.18), GOLD)
text(s, Inches(0.85), Inches(0.7), Inches(11.6), Inches(1),
     [[("KEY TAKEAWAYS", 14, GOLD, True)], [("Five Things to Remember", 34, WHITE, True)]],
     space=3.0)
takeaways = [
    ("Start now", "Time in the market is your biggest advantage."),
    ("Stay diversified", "Own many things, never bet it all on one."),
    ("Keep it cheap & simple", "Low-cost index funds beat complexity."),
    ("Be consistent", "Invest automatically, every single month."),
    ("Ignore the noise", "Stay calm, think in decades, don't panic-sell."),
]
y = Inches(2.2)
for i, (t, b) in enumerate(takeaways):
    rect(s, Inches(0.85), y, Inches(0.55), Inches(0.55), GOLD)
    text(s, Inches(0.85), y, Inches(0.55), Inches(0.55),
         [[(str(i + 1), 20, NAVY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.65), y - Inches(0.05), Inches(10.8), Inches(0.8),
         [[(t, 19, WHITE, True), ("   " + b, 15, GREY, False)]])
    y = Inches(y.inches + 0.86)
text(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.5),
     [[("Disclaimer: Educational content only — not personalized financial advice.",
        12, GREY, False)]], align=PP_ALIGN.CENTER)

out = r"C:\Users\Admin\OneDrive\Documents\max\data\Stock_Market_Investing_Guide.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("SAVED:", out)
print("IMAGES:", list(local.keys()))
